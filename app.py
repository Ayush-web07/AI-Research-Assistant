"""
AI Research Assistant — single-file build.

This file combines the project's 14 modules (database, rag_engine, llm_client,
document_loaders, pdf_vision, image_analysis, quiz_generator, mindmap_generator,
doc_comparison, deep_research, translator, voice_client, export_utils, app) into
one file, so there is nothing to import between local modules and therefore no
way for cross-file copy/paste mistakes to cause circular-import or
missing-attribute errors.

Run with:
    streamlit run app_single_file.py
"""

import asyncio
import base64
import csv
import hashlib
import io
import json
import os
import pickle
import shutil
import sqlite3
import sys
import time
import types
import uuid
import warnings
import wave
from contextlib import contextmanager
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import List

if sys.platform == "win32":
    asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())

os.environ.setdefault("USE_TF", "0")
os.environ.setdefault("TRANSFORMERS_NO_ADVISORY_WARNINGS", "1")
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("TF_CPP_MIN_LOG_LEVEL", "3")
os.environ.setdefault("TF_ENABLE_ONEDNN_OPTS", "0")
warnings.filterwarnings("ignore", category=FutureWarning)

import re

import requests
import streamlit as st
from dotenv import load_dotenv, set_key

import faiss
import numpy as np
import pdfplumber
from sentence_transformers import SentenceTransformer

import fitz  # PyMuPDF
from PIL import Image

import docx  # python-docx
import openpyxl
from pptx import Presentation  # python-pptx

from json_repair import repair_json

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable


# ============================================================================
# --- from database.py =================================================
# ============================================================================

import sqlite3
import json
from datetime import datetime
from pathlib import Path
from contextlib import contextmanager

DB_PATH = Path(__file__).parent / "data" / "research_assistant.db"


@contextmanager
def get_connection():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
        conn.commit()
    finally:
        conn.close()


def init_db():
    """Create tables if they don't already exist. Safe to call every app start."""
    DB_PATH.parent.mkdir(parents=True, exist_ok=True)
    with get_connection() as conn:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                filename TEXT NOT NULL,
                num_pages INTEGER,
                num_chunks INTEGER,
                uploaded_at TEXT NOT NULL,
                session_id TEXT NOT NULL
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT NOT NULL,
                question TEXT NOT NULL,
                answer TEXT NOT NULL,
                sources_json TEXT,
                created_at TEXT NOT NULL
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS tool_usage (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT NOT NULL,
                tool_key TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
        """)


def add_document(filename: str, num_pages: int, num_chunks: int, session_id: str) -> int:
    with get_connection() as conn:
        cur = conn.execute(
            "INSERT INTO documents (filename, num_pages, num_chunks, uploaded_at, session_id) "
            "VALUES (?, ?, ?, ?, ?)",
            (filename, num_pages, num_chunks, datetime.utcnow().isoformat(), session_id),
        )
        return cur.lastrowid


def get_documents(session_id: str):
    with get_connection() as conn:
        rows = conn.execute(
            "SELECT * FROM documents WHERE session_id = ? ORDER BY uploaded_at DESC",
            (session_id,),
        ).fetchall()
        return [dict(r) for r in rows]


def add_chat_entry(session_id: str, question: str, answer: str, sources: list) -> int:
    with get_connection() as conn:
        cur = conn.execute(
            "INSERT INTO chat_history (session_id, question, answer, sources_json, created_at) "
            "VALUES (?, ?, ?, ?, ?)",
            (session_id, question, answer, json.dumps(sources), datetime.utcnow().isoformat()),
        )
        return cur.lastrowid


def get_chat_history(session_id: str):
    with get_connection() as conn:
        rows = conn.execute(
            "SELECT * FROM chat_history WHERE session_id = ? ORDER BY created_at ASC",
            (session_id,),
        ).fetchall()
        result = []
        for r in rows:
            d = dict(r)
            d["sources"] = json.loads(d.pop("sources_json") or "[]")
            result.append(d)
        return result


def clear_chat_history(session_id: str):
    """Clear only the chat conversation (used by the per-tool 'Clear chat' button
    in Chat with PDF / Voice Chat) — leaves uploaded documents and other tool
    results untouched, unlike clear_session() which wipes everything."""
    with get_connection() as conn:
        conn.execute("DELETE FROM chat_history WHERE session_id = ?", (session_id,))


def clear_session(session_id: str):
    with get_connection() as conn:
        conn.execute("DELETE FROM chat_history WHERE session_id = ?", (session_id,))
        conn.execute("DELETE FROM documents WHERE session_id = ?", (session_id,))
        conn.execute("DELETE FROM tool_usage WHERE session_id = ?", (session_id,))


def log_tool_usage(session_id: str, tool_key: str) -> int:
    """Record one successful use of a tool (e.g. 'quiz', 'mindmap', 'compare',
    'research', 'image'). Called by app.py right after a tool completes without
    error — this only counts real, successful runs, not button clicks that failed."""
    with get_connection() as conn:
        cur = conn.execute(
            "INSERT INTO tool_usage (session_id, tool_key, created_at) VALUES (?, ?, ?)",
            (session_id, tool_key, datetime.utcnow().isoformat()),
        )
        return cur.lastrowid


def get_tool_usage_counts(session_id: str) -> dict:
    """Returns {tool_key: count} for the given session. Keys only exist for
    tools that have been used at least once."""
    with get_connection() as conn:
        rows = conn.execute(
            "SELECT tool_key, COUNT(*) as cnt FROM tool_usage WHERE session_id = ? GROUP BY tool_key",
            (session_id,),
        ).fetchall()
        return {r["tool_key"]: r["cnt"] for r in rows}

# ============================================================================
# --- from pdf_vision.py ===============================================
# ============================================================================

import base64
import io

import fitz  # PyMuPDF
from PIL import Image

MAX_IMAGE_DIMENSION = 1024  # px, caps vision API token cost


class PDFVisionError(Exception):
    pass


def resize_image_bytes(image_bytes: bytes, max_dimension: int = MAX_IMAGE_DIMENSION) -> bytes:
    """Downscale an image so its longest side is at most max_dimension pixels,
    re-encoded as PNG. This is applied right before sending any image to the
    vision API — whether it's a rendered PDF page, an extracted embedded image,
    or a raw uploaded photo — so token cost stays predictable regardless of
    source resolution. Charts and text stay readable at this size; a 1024px
    chart is still perfectly legible to a vision model."""
    img = Image.open(io.BytesIO(image_bytes))
    if img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGB")

    width, height = img.size
    longest_side = max(width, height)
    if longest_side > max_dimension:
        scale = max_dimension / longest_side
        new_size = (max(1, int(width * scale)), max(1, int(height * scale)))
        img = img.resize(new_size, Image.LANCZOS)

    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()


def get_page_count(pdf_path: str) -> int:
    with fitz.open(pdf_path) as doc:
        return doc.page_count


def render_page_png(pdf_path: str, page_number: int, zoom: float = 2.0) -> bytes:
    """Render a single page (1-indexed) to PNG bytes at the given zoom factor.
    zoom=2.0 is roughly 144 DPI — a good balance of clarity vs file size for
    vision models (Groq's limit is 20MB per image; a rendered page at this
    zoom is typically well under 1MB)."""
    with fitz.open(pdf_path) as doc:
        if not (1 <= page_number <= doc.page_count):
            raise PDFVisionError(f"Page {page_number} is out of range (1-{doc.page_count}).")
        page = doc.load_page(page_number - 1)
        matrix = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=matrix)
        return pix.tobytes("png")


def extract_embedded_images(pdf_path: str, page_number: int) -> list:
    """Extract raw embedded raster images from a single page (1-indexed), each
    normalized to PNG bytes. Images that fail to extract cleanly (unusual color
    spaces, corrupt streams, etc.) are silently skipped rather than raising,
    since a single bad image shouldn't block extraction of the rest."""
    images = []
    with fitz.open(pdf_path) as doc:
        if not (1 <= page_number <= doc.page_count):
            raise PDFVisionError(f"Page {page_number} is out of range (1-{doc.page_count}).")
        page = doc.load_page(page_number - 1)
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image.get("ext", "png").lower()
                if ext not in ("png", "jpg", "jpeg"):
                    # normalize unusual formats (e.g. CMYK, unusual color spaces)
                    # to PNG by re-rendering through a Pixmap
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n - pix.alpha >= 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    image_bytes = pix.tobytes("png")
                images.append(image_bytes)
            except Exception:
                continue
    return images


def to_base64(image_bytes: bytes) -> str:
    return base64.b64encode(image_bytes).decode("utf-8")

# ============================================================================
# --- from llm_client.py ===============================================
# ============================================================================

import json
import os
import re
import time
import requests

GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL = "openai/gpt-oss-20b"
GROQ_VISION_MODEL = "qwen/qwen3.6-27b"  # only Groq model currently used that accepts image input

HF_API_URL = "https://api-inference.huggingface.co/models/mistralai/Mistral-7B-Instruct-v0.3"

# Groq's free tier has tight per-minute token limits, and vision calls in particular
# are token-heavy (the image itself counts). Automatically retry once or twice on a
# 429, using the exact wait time Groq tells us in the error message, before giving up.
GROQ_MAX_RETRIES = 2
GROQ_RETRY_JITTER = 0.5  # small buffer added on top of Groq's suggested wait

SYSTEM_PROMPT_LLMCLIENT = (
    "You are a precise research assistant. Answer the user's question using ONLY the "
    "provided document excerpts. If the excerpts don't contain the answer, say so clearly "
    "instead of guessing. When you use information from an excerpt, mention which source "
    "number it came from, e.g. [Source 1]. Be thorough and complete: if the user asks for "
    "a specific number of items (e.g. '20 questions', 'list all the...'), provide exactly "
    "that many, fully numbered, without stopping early or summarizing partway through. "
    "Otherwise, keep answers well-structured and free of unnecessary padding."
)


class LLMError(Exception):
    pass


def _parse_retry_after_seconds(error_text: str, default: float = 2.0) -> float:
    """Groq's 429 responses include a hint like 'Please try again in 3.375s.'
    Parse that out so we wait exactly as long as needed, not an arbitrary guess."""
    match = re.search(r"try again in ([\d.]+)s", error_text)
    if match:
        try:
            return float(match.group(1)) + GROQ_RETRY_JITTER
        except ValueError:
            pass
    return default


def _groq_post(payload: dict, api_key: str, url: str = None, stream: bool = False, timeout: int = 90):
    """POST JSON to a Groq endpoint with automatic retry on 429 (rate limit)
    responses. Defaults to the chat completions endpoint; pass `url` for other
    Groq endpoints (e.g. text-to-speech). Non-429 errors (and the final attempt
    after exhausting retries) are returned as-is for the caller to handle/raise."""
    url = url or GROQ_API_URL
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    attempt = 0
    while True:
        resp = requests.post(url, headers=headers, json=payload, timeout=timeout, stream=stream)
        if resp.status_code == 429 and attempt < GROQ_MAX_RETRIES:
            wait_seconds = _parse_retry_after_seconds(resp.text)
            time.sleep(wait_seconds)
            attempt += 1
            continue
        return resp


def _groq_post_multipart(url: str, api_key: str, files: dict, data: dict, timeout: int = 90):
    """POST multipart/form-data to a Groq endpoint (used for audio file uploads,
    e.g. speech-to-text) with the same 429 retry behavior as _groq_post."""
    headers = {"Authorization": f"Bearer {api_key}"}
    attempt = 0
    while True:
        resp = requests.post(url, headers=headers, files=files, data=data, timeout=timeout)
        if resp.status_code == 429 and attempt < GROQ_MAX_RETRIES:
            wait_seconds = _parse_retry_after_seconds(resp.text)
            time.sleep(wait_seconds)
            attempt += 1
            continue
        return resp


def _build_context(chunks) -> str:
    """Format retrieved chunks into a numbered context block the model can cite."""
    parts = []
    for i, c in enumerate(chunks, start=1):
        parts.append(f"[Source {i}] ({c.filename}, page {c.page}):\n{c.text}")
    return "\n\n".join(parts)


def _call_groq(question: str, context: str, api_key: str) -> str:
    payload = {
        "model": GROQ_MODEL,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT_LLMCLIENT},
            {"role": "user", "content": f"Document excerpts:\n\n{context}\n\nQuestion: {question}"},
        ],
        "temperature": 0.2,
        "max_tokens": 3000,
    }
    resp = _groq_post(payload, api_key, timeout=60)
    if resp.status_code != 200:
        raise LLMError(f"Groq API error {resp.status_code}: {resp.text[:300]}")
    data = resp.json()
    return data["choices"][0]["message"]["content"].strip()


def _call_hf(question: str, context: str, api_key: str) -> str:
    headers = {"Authorization": f"Bearer {api_key}"}
    prompt = (
        f"<s>[INST] {SYSTEM_PROMPT_LLMCLIENT}\n\nDocument excerpts:\n\n{context}\n\n"
        f"Question: {question} [/INST]"
    )
    payload = {"inputs": prompt, "parameters": {"max_new_tokens": 1500, "temperature": 0.2}}
    resp = requests.post(HF_API_URL, headers=headers, json=payload, timeout=60)
    if resp.status_code != 200:
        raise LLMError(f"HF API error {resp.status_code}: {resp.text[:300]}")
    data = resp.json()
    if isinstance(data, list) and data and "generated_text" in data[0]:
        text = data[0]["generated_text"]
        return text.split("[/INST]")[-1].strip()
    raise LLMError(f"Unexpected HF response format: {data}")


def _call_groq_raw(system_prompt: str, user_prompt: str, api_key: str, max_tokens: int = 3000) -> str:
    """Like _call_groq but takes a fully custom system/user prompt instead of the
    citation-style Q&A format. Used by features that need structured output
    (e.g. JSON) rather than a grounded, cited answer."""
    payload = {
        "model": GROQ_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "temperature": 0.3,
        "max_tokens": max_tokens,
    }
    resp = _groq_post(payload, api_key, timeout=90)
    if resp.status_code != 200:
        raise LLMError(f"Groq API error {resp.status_code}: {resp.text[:300]}")
    data = resp.json()
    return data["choices"][0]["message"]["content"].strip()


def _call_hf_raw(system_prompt: str, user_prompt: str, api_key: str, max_new_tokens: int = 1500) -> str:
    """HF equivalent of _call_groq_raw."""
    headers = {"Authorization": f"Bearer {api_key}"}
    prompt = f"<s>[INST] {system_prompt}\n\n{user_prompt} [/INST]"
    payload = {"inputs": prompt, "parameters": {"max_new_tokens": max_new_tokens, "temperature": 0.3}}
    resp = requests.post(HF_API_URL, headers=headers, json=payload, timeout=90)
    if resp.status_code != 200:
        raise LLMError(f"HF API error {resp.status_code}: {resp.text[:300]}")
    data = resp.json()
    if isinstance(data, list) and data and "generated_text" in data[0]:
        text = data[0]["generated_text"]
        return text.split("[/INST]")[-1].strip()
    raise LLMError(f"Unexpected HF response format: {data}")


def generate_raw(system_prompt: str, user_prompt: str, max_tokens: int = 3000) -> str:
    """Generic completion call with a fully custom system/user prompt (no citation
    formatting, no retrieved-chunks assumption). Same Groq -> HF provider fallback
    as generate_answer. Used for structured-output features like quiz generation.
    """
    provider = os.environ.get("LLM_PROVIDER", "auto").lower()
    groq_key = os.environ.get("GROQ_API_KEY")
    hf_key = os.environ.get("HF_API_TOKEN")

    if provider == "groq":
        if not groq_key:
            raise LLMError("LLM_PROVIDER=groq but GROQ_API_KEY is not set.")
        return _call_groq_raw(system_prompt, user_prompt, groq_key, max_tokens=max_tokens)

    if provider == "hf":
        if not hf_key:
            raise LLMError("LLM_PROVIDER=hf but HF_API_TOKEN is not set.")
        return _call_hf_raw(system_prompt, user_prompt, hf_key, max_new_tokens=min(max_tokens, 1500))

    errors = []
    if groq_key:
        try:
            return _call_groq_raw(system_prompt, user_prompt, groq_key, max_tokens=max_tokens)
        except Exception as e:
            errors.append(f"Groq: {e}")
    if hf_key:
        try:
            return _call_hf_raw(system_prompt, user_prompt, hf_key, max_new_tokens=min(max_tokens, 1500))
        except Exception as e:
            errors.append(f"HF: {e}")

    if not groq_key and not hf_key and not errors:
        raise LLMError(
            "No LLM provider available. Set GROQ_API_KEY (free at console.groq.com) "
            "or HF_API_TOKEN (free at huggingface.co/settings/tokens)."
        )
    raise LLMError("All configured LLM providers failed: " + " | ".join(errors))


def _call_groq_stream(question: str, context: str, api_key: str):
    """Stream tokens from Groq as they're generated (Server-Sent Events)."""
    payload = {
        "model": GROQ_MODEL,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT_LLMCLIENT},
            {"role": "user", "content": f"Document excerpts:\n\n{context}\n\nQuestion: {question}"},
        ],
        "temperature": 0.2,
        "max_tokens": 3000,
        "stream": True,
    }
    resp = _groq_post(payload, api_key, stream=True, timeout=60)
    if resp.status_code != 200:
        raise LLMError(f"Groq API error {resp.status_code}: {resp.text[:300]}")

    for raw_line in resp.iter_lines():
        if not raw_line:
            continue
        line = raw_line.decode("utf-8") if isinstance(raw_line, bytes) else raw_line
        if not line.startswith("data:"):
            continue
        data_str = line[len("data:"):].strip()
        if data_str == "[DONE]":
            break
        try:
            event = json.loads(data_str)
        except json.JSONDecodeError:
            continue
        delta = event.get("choices", [{}])[0].get("delta", {})
        content = delta.get("content")
        if content:
            yield content


def generate_answer_stream(question: str, chunks):
    """Generator version of generate_answer — yields answer text piece by piece as it's
    generated, so the UI can render it word-by-word instead of waiting for the full answer.

    Provider selection mirrors generate_answer's LLM_PROVIDER logic.

    Hugging Face's Inference API doesn't support reliable token streaming for arbitrary
    models, so that path yields the full answer as a single chunk instead of failing.
    """
    if not chunks:
        yield "I couldn't find any relevant content in the uploaded documents to answer that."
        return

    context = _build_context(chunks)
    provider = os.environ.get("LLM_PROVIDER", "auto").lower()

    groq_key = os.environ.get("GROQ_API_KEY")
    hf_key = os.environ.get("HF_API_TOKEN")

    if provider == "groq":
        if not groq_key:
            raise LLMError("LLM_PROVIDER=groq but GROQ_API_KEY is not set.")
        yield from _call_groq_stream(question, context, groq_key)
        return

    if provider == "hf":
        if not hf_key:
            raise LLMError("LLM_PROVIDER=hf but HF_API_TOKEN is not set.")
        yield _call_hf(question, context, hf_key)
        return

    # auto mode: pick the first available provider, same priority as generate_answer
    if groq_key:
        yield from _call_groq_stream(question, context, groq_key)
        return
    if hf_key:
        yield _call_hf(question, context, hf_key)
        return

    raise LLMError(
        "No LLM provider available. Set GROQ_API_KEY (free at console.groq.com) "
        "or HF_API_TOKEN (free at huggingface.co/settings/tokens)."
    )


def generate_answer(question: str, chunks) -> str:
    """Generate an answer grounded in retrieved chunks.

    Provider selection follows LLM_PROVIDER (see module docstring). In "auto" mode
    (the default) it tries Groq, then Hugging Face, using whichever is configured first.
    """
    if not chunks:
        return "I couldn't find any relevant content in the uploaded documents to answer that."

    context = _build_context(chunks)
    provider = os.environ.get("LLM_PROVIDER", "auto").lower()

    groq_key = os.environ.get("GROQ_API_KEY")
    hf_key = os.environ.get("HF_API_TOKEN")

    if provider == "groq":
        if not groq_key:
            raise LLMError("LLM_PROVIDER=groq but GROQ_API_KEY is not set.")
        return _call_groq(question, context, groq_key)

    if provider == "hf":
        if not hf_key:
            raise LLMError("LLM_PROVIDER=hf but HF_API_TOKEN is not set.")
        return _call_hf(question, context, hf_key)

    # provider == "auto": try each in order, collecting errors as we go
    errors = []
    if groq_key:
        try:
            return _call_groq(question, context, groq_key)
        except Exception as e:
            errors.append(f"Groq: {e}")
    if hf_key:
        try:
            return _call_hf(question, context, hf_key)
        except Exception as e:
            errors.append(f"HF: {e}")

    if not groq_key and not hf_key and not errors:
        raise LLMError(
            "No LLM provider available. Set GROQ_API_KEY (free at console.groq.com) "
            "or HF_API_TOKEN (free at huggingface.co/settings/tokens)."
        )
    raise LLMError("All configured LLM providers failed: " + " | ".join(errors))


def summarize_document(full_text: str) -> str:
    """Summarize an entire document (used for the 'AI summarizes research papers' feature).
    Truncates very long documents to stay within context limits of free-tier models."""
    max_chars = 12000
    truncated = full_text[:max_chars]
    question = (
        "Summarize this research paper in 5-8 bullet points, covering: the research "
        "question, methodology, key findings, and conclusions."
    )
    class _FakeChunk:
        def __init__(self, text):
            self.text = text
            self.filename = "document"
            self.page = 1
    return generate_answer(question, [_FakeChunk(truncated)])


def _call_groq_vision(image_b64: str, mime_type: str, prompt: str, api_key: str, max_tokens: int = 1500) -> str:
    """Call Groq's vision-capable model with a base64-encoded image plus a text prompt."""
    payload = {
        "model": GROQ_VISION_MODEL,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_b64}"}},
                ],
            }
        ],
        "temperature": 0.3,
        "max_tokens": max_tokens,
    }
    resp = _groq_post(payload, api_key, timeout=90)
    if resp.status_code != 200:
        raise LLMError(f"Groq Vision API error {resp.status_code}: {resp.text[:300]}")
    data = resp.json()
    return data["choices"][0]["message"]["content"].strip()


def analyze_image(image_b64: str, mime_type: str = "image/png", prompt: str = None, max_tokens: int = 1500) -> str:
    """Analyze an image (e.g. a rendered PDF page or an extracted figure) using a
    vision-capable model. Only Groq is used here — Hugging Face's free Inference
    API doesn't reliably support vision across arbitrary models, so there's no
    fallback provider for this specific feature.

    Args:
        image_b64: base64-encoded image bytes (no data: prefix).
        mime_type: e.g. "image/png" or "image/jpeg".
        prompt: what to ask about the image. Defaults to a general description
            + chart/table data extraction prompt if not provided.
        max_tokens: response length cap.

    Raises:
        LLMError: GROQ_API_KEY isn't set, or the API call failed.
    """
    groq_key = os.environ.get("GROQ_API_KEY")
    if not groq_key:
        raise LLMError(
            "Image analysis requires a Groq API key (GROQ_API_KEY) — Groq is currently "
            "the only configured provider with vision support in this app."
        )
    default_prompt = (
        "Describe what's in this image in detail. If it contains a chart, graph, or "
        "table, extract and explain the key data, trends, axis labels, and any "
        "numeric values visible. If it's a diagram, explain what it illustrates. "
        "If it's a photo or figure, describe its content and relevance."
    )
    return _call_groq_vision(image_b64, mime_type, prompt or default_prompt, groq_key, max_tokens=max_tokens)

# ============================================================================
# --- from document_loaders.py =========================================
# ============================================================================

import csv
import io
from pathlib import Path

import docx  # python-docx
import openpyxl
from pptx import Presentation  # python-pptx


# Extensions handled as plain UTF-8 text (source code, markup, config, etc.)
PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".go", ".rb", ".php", ".rs", ".swift", ".kt", ".sql",
    ".html", ".css", ".json", ".yaml", ".yml", ".xml", ".sh", ".bat",
}

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".tif"}

# All extensions this module claims to support, for use in file_uploader's `type=` list.
# Note: ".pdf" is included here for that purpose, but load_document() below does NOT
# handle it — PDF extraction stays in rag_engine.py (via pdfplumber, already tested)
# to avoid a circular import (rag_engine imports this module, not the other way around).
# VectorStore.add_document() special-cases ".pdf" before falling through to this module.
SUPPORTED_EXTENSIONS = (
    {".pdf", ".docx", ".pptx", ".xlsx", ".csv"} | PLAIN_TEXT_EXTENSIONS | IMAGE_EXTENSIONS
)


class DocumentLoadError(Exception):
    """Raised when a file can't be parsed into text content."""
    pass


def _load_docx(path: str) -> list:
    document = docx.Document(path)
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    # Tables often carry real content (e.g. resumes, spec sheets) — include them too
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    full_text = "\n".join(paragraphs)
    return [full_text] if full_text.strip() else []


def _load_pptx(path: str) -> list:
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        pieces = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text.strip():
                    pieces.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        pieces.append(row_text)
        slides_text.append("\n".join(pieces))
    return slides_text


def _load_xlsx(path: str) -> list:
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    sheets_text = []
    for sheet in wb.worksheets:
        lines = [f"Sheet: {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
        sheets_text.append("\n".join(lines))
    wb.close()
    return sheets_text


def _load_csv(path: str) -> list:
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        lines = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
    return ["\n".join(lines)] if lines else []


def _load_plain_text(path: str) -> list:
    with open(path, encoding="utf-8", errors="replace") as f:
        content = f.read()
    return [content] if content.strip() else []


def _load_image(path: str, ext: str) -> list:
    """Describe an uploaded image using the vision model so its content becomes
    searchable text. Raises LLMError if no Groq key is configured (the only
    vision-capable provider in this app)."""
    with open(path, "rb") as f:
        raw_bytes = f.read()
    import base64
    # resize_image_bytes always re-encodes to PNG, so the MIME type sent to the
    # API must be image/png regardless of the original upload's extension
    raw_bytes = resize_image_bytes(raw_bytes)
    b64 = base64.b64encode(raw_bytes).decode("utf-8")
    mime_type = "image/png"

    prompt = (
        "Describe this image thoroughly for a search index: transcribe any visible text "
        "verbatim (OCR), describe any charts/graphs/tables with their actual data, and "
        "describe any diagrams, photos, or handwriting. Be comprehensive — this is the "
        "only representation of this image that will be searchable."
    )
    description = analyze_image(b64, mime_type=mime_type, prompt=prompt, max_tokens=1500)
    return [description]


def load_document(file_path: str, filename: str) -> list:
    """Load a file and return a list of page/section text strings.

    Dispatches based on file extension. Does NOT handle ".pdf" — that's handled
    directly by rag_engine.py's VectorStore.add_document(), which special-cases
    it before falling through here (see the SUPPORTED_EXTENSIONS comment above).

    Raises:
        DocumentLoadError: unsupported extension (including ".pdf", if called
            directly with one — callers should route PDFs elsewhere).
        LLMError: an image needs analysis but no vision provider is configured.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        raise DocumentLoadError(
            "load_document() doesn't handle PDFs directly — use "
            "VectorStore.add_document(), which routes PDFs to rag_engine's own "
            "pdfplumber-based extraction."
        )
    if ext == ".docx":
        return _load_docx(file_path)
    if ext == ".pptx":
        return _load_pptx(file_path)
    if ext == ".xlsx":
        return _load_xlsx(file_path)
    if ext == ".csv":
        return _load_csv(file_path)
    if ext in PLAIN_TEXT_EXTENSIONS:
        return _load_plain_text(file_path)
    if ext in IMAGE_EXTENSIONS:
        return _load_image(file_path, ext)

    raise DocumentLoadError(
        f"Unsupported file type: '{ext}'. Supported: PDF, DOCX, PPTX, XLSX, CSV, "
        f"images (jpg/png/webp/bmp/tiff), and plain text/code files."
    )

# ============================================================================
# --- from rag_engine.py ===============================================
# ============================================================================

import json
import os
import warnings
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import List

# This project only needs the PyTorch backend for embeddings. Some environments
# also have TensorFlow installed for unrelated work, which makes `transformers`
# (a dependency of sentence-transformers) auto-detect and try to load it —
# producing noisy oneDNN/Keras warnings that have nothing to do with this app.
# Telling it up front to skip TF avoids that whole chain, and is safe because
# we never use TF here.
os.environ.setdefault("USE_TF", "0")
os.environ.setdefault("TRANSFORMERS_NO_ADVISORY_WARNINGS", "1")
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("TF_CPP_MIN_LOG_LEVEL", "3")
os.environ.setdefault("TF_ENABLE_ONEDNN_OPTS", "0")
warnings.filterwarnings("ignore", category=FutureWarning)

import faiss
import numpy as np
import pdfplumber
from sentence_transformers import SentenceTransformer

EMBED_MODEL_NAME = "all-MiniLM-L6-v2"  # small, fast, free, runs locally
CHUNK_SIZE = 800        # characters per chunk
CHUNK_OVERLAP = 150     # overlap so sentences at boundaries aren't lost
INDEX_DIR = Path(__file__).parent / "data" / "indexes"


@dataclass
class Chunk:
    text: str
    filename: str
    page: int
    chunk_id: int


@dataclass
class RetrievedChunk:
    text: str
    filename: str
    page: int
    score: float


class EmbeddingModel:
    """Lazy singleton wrapper so the model is only loaded once per process."""
    _instance = None

    @classmethod
    def get(cls) -> SentenceTransformer:
        if cls._instance is None:
            cls._instance = SentenceTransformer(EMBED_MODEL_NAME)
        return cls._instance


def extract_pages(pdf_path: str) -> List[str]:
    """Return a list of raw text strings, one per page."""
    pages = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages.append(text)
    return pages


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Simple sliding-window chunker with overlap. Splits on whitespace boundaries
    where possible to avoid cutting words in half."""
    text = " ".join(text.split())  # normalize whitespace
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            # try to break on a space near the boundary
            space_idx = text.rfind(" ", start, end)
            if space_idx > start:
                end = space_idx
        chunks.append(text[start:end].strip())
        start = end - overlap if end - overlap > start else end
    return [c for c in chunks if c]


def build_chunks_from_pdf(pdf_path: str, filename: str) -> List[Chunk]:
    pages = extract_pages(pdf_path)
    chunks = []
    chunk_id = 0
    for page_num, page_text in enumerate(pages, start=1):
        for piece in chunk_text(page_text):
            chunks.append(Chunk(text=piece, filename=filename, page=page_num, chunk_id=chunk_id))
            chunk_id += 1
    return chunks, len(pages)


class VectorStore:
    """FAISS-backed vector store for one user session. Holds chunks from
    potentially multiple uploaded PDFs so questions can search across all of them."""

    def __init__(self, session_id: str):
        self.session_id = session_id
        self.index_path = INDEX_DIR / f"{session_id}.faiss"
        self.meta_path = INDEX_DIR / f"{session_id}.json"
        INDEX_DIR.mkdir(parents=True, exist_ok=True)

        model = EmbeddingModel.get()
        if hasattr(model, "get_embedding_dimension"):
            self.dim = model.get_embedding_dimension()
        else:
            # older sentence-transformers versions don't have the renamed method
            self.dim = model.get_sentence_embedding_dimension()
        self.chunks: List[Chunk] = []
        self.index = faiss.IndexFlatIP(self.dim)  # cosine similarity via normalized vectors

        self._load_if_exists()

    def _load_if_exists(self):
        if self.index_path.exists() and self.meta_path.exists():
            self.index = faiss.read_index(str(self.index_path))
            with open(self.meta_path, "r", encoding="utf-8") as f:
                raw = json.load(f)
            self.chunks = [Chunk(**d) for d in raw]

    def _save(self):
        faiss.write_index(self.index, str(self.index_path))
        # JSON, not pickle: pickle identifies classes by module+qualname and
        # requires the *exact same class object* at load time. That's fragile
        # here because Streamlit re-executes the whole script (redefining
        # Chunk fresh) on every rerun in the single-file build — pickling a
        # Chunk from one rerun and loading it in the next raises
        # "it's not the same object" even though the class is identical.
        # JSON has no such notion of class identity, so this sidesteps the
        # issue entirely regardless of how the app is structured.
        with open(self.meta_path, "w", encoding="utf-8") as f:
            json.dump([asdict(c) for c in self.chunks], f)

    def _add_chunks(self, new_chunks: List[Chunk]):
        """Shared embed-and-index logic used by both add_pdf and add_document."""
        if not new_chunks:
            return
        texts = [c.text for c in new_chunks]
        embeddings = EmbeddingModel.get().encode(texts, normalize_embeddings=True, show_progress_bar=False)
        embeddings = np.array(embeddings, dtype="float32")

        self.index.add(embeddings)
        self.chunks.extend(new_chunks)
        self._save()

    def add_pdf(self, pdf_path: str, filename: str):
        new_chunks, num_pages = build_chunks_from_pdf(pdf_path, filename)
        self._add_chunks(new_chunks)
        return len(new_chunks), num_pages

    def add_document(self, file_path: str, filename: str):
        """Generic ingestion for any supported file type (PDF or otherwise).
        Returns (num_chunks, num_pages) just like add_pdf, so callers don't need
        to branch on file type.

        Raises:
            document_loaders.DocumentLoadError: unsupported extension
            LLMError: an image needs vision analysis but no provider is configured
        """
        ext = Path(filename).suffix.lower()
        if ext == ".pdf":
            return self.add_pdf(file_path, filename)

        # local import avoids a circular import (document_loaders -> llm_client,
        # nothing imports rag_engine, so this is safe to import lazily here)
        pages = load_document(file_path, filename)

        chunks = []
        chunk_id = 0
        for page_num, page_text in enumerate(pages, start=1):
            for piece in chunk_text(page_text):
                chunks.append(Chunk(text=piece, filename=filename, page=page_num, chunk_id=chunk_id))
                chunk_id += 1

        self._add_chunks(chunks)
        return len(chunks), len(pages)

    def search(self, query: str, top_k: int = 5) -> List[RetrievedChunk]:
        if self.index.ntotal == 0:
            return []
        query_vec = EmbeddingModel.get().encode([query], normalize_embeddings=True)
        query_vec = np.array(query_vec, dtype="float32")
        scores, indices = self.index.search(query_vec, min(top_k, self.index.ntotal))

        results = []
        for score, idx in zip(scores[0], indices[0]):
            if idx == -1:
                continue
            c = self.chunks[idx]
            results.append(RetrievedChunk(text=c.text, filename=c.filename, page=c.page, score=float(score)))
        return results

    def has_documents(self) -> bool:
        return self.index.ntotal > 0

    def list_filenames(self) -> List[str]:
        return sorted(set(c.filename for c in self.chunks))

    def clear(self):
        self.index = faiss.IndexFlatIP(self.dim)
        self.chunks = []
        for p in (self.index_path, self.meta_path):
            if p.exists():
                p.unlink()

# ============================================================================
# --- from image_analysis.py ===========================================
# ============================================================================

def describe_page(pdf_path: str, page_number: int, question: str = None, zoom: float = 2.0) -> str:
    """Render the given page as an image and describe it (or answer a specific
    question about it) using the vision model. This is the primary path — it
    works for both real embedded images AND vector-drawn charts/diagrams that
    have no separate embedded image object.

    Raises:
        PDFVisionError: invalid page number
        LLMError: no Groq API key configured, or the vision call failed
    """
    image_bytes = render_page_png(pdf_path, page_number, zoom=zoom)
    image_bytes = resize_image_bytes(image_bytes)  # cap resolution -> cap token cost
    b64 = to_base64(image_bytes)
    prompt = None
    if question and question.strip():
        prompt = (
            f"Looking at this page from a document, answer the following question. "
            f"If the page contains a chart, graph, or table, use the actual data shown "
            f"to answer precisely, quoting specific numbers/labels where relevant.\n\n"
            f"Question: {question.strip()}"
        )
    return analyze_image(b64, mime_type="image/png", prompt=prompt)


def describe_embedded_images(pdf_path: str, page_number: int, question: str = None) -> list:
    """Extract and describe each individually embedded raster image on a page.

    Returns a list of description strings, one per image found. Returns an
    empty list if the page has no embedded raster images (common for pages
    where charts/diagrams are vector-drawn rather than embedded pictures —
    use describe_page() instead in that case, which handles both).

    Individual image analysis failures (e.g. a transient API error on one of
    several images) are captured as an error string in that image's slot
    rather than aborting the whole batch.
    """
    images = extract_embedded_images(pdf_path, page_number)
    descriptions = []
    prompt = None
    if question and question.strip():
        prompt = f"Looking at this image, answer: {question.strip()}"

    for img_bytes in images:
        img_bytes = resize_image_bytes(img_bytes)  # cap resolution -> cap token cost
        b64 = to_base64(img_bytes)
        try:
            desc = analyze_image(b64, mime_type="image/png", prompt=prompt)
        except LLMError as e:
            desc = f"(Could not analyze this image: {e})"
        descriptions.append(desc)

    return descriptions

# ============================================================================
# --- from quiz_generator.py ===========================================
# ============================================================================

import json
import re

from json_repair import repair_json


QUIZ_TYPES = ["MCQ", "Coding Questions", "Interview Questions", "Flashcards"]
DIFFICULTIES = ["Easy", "Medium", "Hard"]

MAX_SOURCE_CHARS_QUIZ = 8000  # keep prompt within free-tier context limits (see token math in generate_quiz)

SYSTEM_PROMPT_QUIZ = (
    "You are a quiz-generation assistant. You output ONLY a valid JSON array, nothing else — "
    "no explanations, no markdown formatting, no code fences, no text before or after the "
    "JSON. If you include anything other than the JSON array, the output is unusable."
)


class QuizGenerationError(Exception):
    """Raised when the model's output can't be parsed into a valid quiz."""
    pass


def _extract_json(text: str):
    """Strip markdown code fences / stray prose and parse the first JSON array found."""
    text = text.strip()

    fence_match = re.search(r"```(?:json)?\s*(.*?)```", text, re.DOTALL)
    if fence_match:
        text = fence_match.group(1).strip()

    start = text.find("[")
    end = text.rfind("]")
    if start == -1 or end == -1 or end < start:
        preview = text[:200].replace("\n", " ") if text else "(empty response)"
        raise QuizGenerationError(
            f"The model didn't respond with JSON — it wrote this instead: \"{preview}"
            f"{'...' if len(text) > 200 else ''}\""
        )
    json_str = text[start:end + 1]
    try:
        return json.loads(json_str, strict=False)
    except json.JSONDecodeError:
        pass  # fall through to repair attempt below

    try:
        return json.loads(repair_json(json_str), strict=False)
    except (json.JSONDecodeError, Exception) as e:
        raise QuizGenerationError(
            f"Couldn't parse the model's output as JSON, even after attempting repair: {e}. "
            f"Try again, or reduce the number of questions requested."
        )


def _build_instructions(quiz_type: str, difficulty: str, num_items: int) -> str:
    if quiz_type == "MCQ":
        schema = (
            '[{"question": "...", "options": ["...", "...", "...", "..."], '
            '"correct_index": 0, "explanation": "..."}]'
        )
        task = (
            f"Generate exactly {num_items} multiple-choice questions at {difficulty} "
            "difficulty testing understanding of the document. Each question must have "
            "exactly 4 options with exactly one correct answer. correct_index is the "
            "0-based index of the correct option."
        )
    elif quiz_type == "Coding Questions":
        schema = '[{"question": "...", "hint": "...", "sample_approach": "..."}]'
        task = (
            f"Generate exactly {num_items} coding or technical exercise questions at "
            f"{difficulty} difficulty, based on the skills, tools, or concepts mentioned "
            "in the document. Include a short hint and a brief sample approach for each."
        )
    elif quiz_type == "Interview Questions":
        schema = '[{"question": "...", "what_to_cover": "..."}]'
        task = (
            f"Generate exactly {num_items} interview questions at {difficulty} difficulty "
            "based on the document. For each, briefly note what a strong answer should cover."
        )
    else:  # Flashcards
        schema = '[{"front": "...", "back": "..."}]'
        task = (
            f"Generate exactly {num_items} flashcards at {difficulty} difficulty summarizing "
            "key facts, terms, or concepts from the document. Keep the front short (a term "
            "or short question) and the back concise (the answer or definition)."
        )

    return (
        f"{task}\n\nRespond with ONLY a valid JSON array matching this schema, no prose "
        f"before or after, no markdown code fences:\n{schema}"
    )


def generate_quiz(
    document_text: str, quiz_type: str, difficulty: str, num_items: int, offset: int = 0
) -> list:
    """Generate a quiz from document_text. Returns a list of dicts whose shape depends
    on quiz_type (see _build_instructions for each schema).

    Args:
        offset: character offset into document_text to start the excerpt window
            from, for documents longer than MAX_SOURCE_CHARS_QUIZ. Used by the UI to
            rotate which part of a large document successive "Generate" clicks
            draw from, so batching toward 50+ items actually covers more of the
            document instead of repeatedly regenerating from the same opening slice.

    Raises:
        ValueError: invalid quiz_type/difficulty/num_items
        LLMError: no provider available or the provider call failed
        QuizGenerationError: the model's output couldn't be parsed as valid JSON
    """
    if quiz_type not in QUIZ_TYPES:
        raise ValueError(f"Unknown quiz type: {quiz_type!r}. Must be one of {QUIZ_TYPES}")
    if difficulty not in DIFFICULTIES:
        raise ValueError(f"Unknown difficulty: {difficulty!r}. Must be one of {DIFFICULTIES}")
    if not (1 <= num_items <= 30):
        raise ValueError("num_items must be between 1 and 30")

    if len(document_text) > MAX_SOURCE_CHARS_QUIZ:
        start = offset % (len(document_text) - MAX_SOURCE_CHARS_QUIZ + 1)
    else:
        start = 0
    truncated = document_text[start:start + MAX_SOURCE_CHARS_QUIZ]
    user_prompt = (
        f"Document content:\n\n{truncated}\n\n"
        f"{_build_instructions(quiz_type, difficulty, num_items)}"
    )

    # Scale the output token budget to what was actually requested, instead of
    # a flat oversized value. A real request for just 10 items on a large
    # document previously still asked for a flat 6000 output tokens, which —
    # combined with input tokens — exceeded Groq's free-tier 8000 TPM ceiling
    # in a single request (a 413 "request too large" error, not a rate limit
    # the automatic 429 retry could help with, since the request itself was
    # simply too big regardless of recent usage).
    tokens_per_item = 180  # generous per-item estimate covering all quiz types
    base_overhead = 300
    max_tokens = min(5000, max(800, num_items * tokens_per_item + base_overhead))

    raw = generate_raw(SYSTEM_PROMPT_QUIZ, user_prompt, max_tokens=max_tokens)
    try:
        items = _extract_json(raw)
    except QuizGenerationError:
        # The model occasionally responds with prose instead of JSON. Retry
        # once with a more insistent reminder before giving up.
        reinforced_prompt = (
            user_prompt + "\n\nREMINDER: your entire response must be ONLY the JSON "
            "array described above — no commentary, no explanation, before or after."
        )
        raw = generate_raw(SYSTEM_PROMPT_QUIZ, reinforced_prompt, max_tokens=max_tokens)
        items = _extract_json(raw)  # let this raise if it fails again

    if not isinstance(items, list) or not items:
        raise QuizGenerationError("The model didn't return a non-empty list of quiz items.")

    # Each item must be a dict (question/options/etc.) — the UI calls .get() on
    # every item, so a non-dict entry (which can happen if json-repair salvaged
    # something from badly malformed output) would otherwise crash the app.
    valid_items = [item for item in items if isinstance(item, dict)]
    if not valid_items:
        raise QuizGenerationError(
            "The model's output couldn't be parsed into usable quiz items. Try again."
        )
    return valid_items

# ============================================================================
# --- from mindmap_generator.py ========================================
# ============================================================================

import json
import re

from json_repair import repair_json


MAX_SOURCE_CHARS_MINDMAP = 8000

SYSTEM_PROMPT_MINDMAP = (
    "You are a mind-map generation assistant. You output ONLY a valid JSON object "
    "representing a hierarchical tree, nothing else — no explanations, no markdown "
    "formatting, no code fences, no text before or after the JSON."
)

SCHEMA_HINT_MINDMAP = (
    '{"topic": "Root Topic", "children": [{"topic": "Subtopic 1", "children": '
    '[{"topic": "Detail A", "children": []}, {"topic": "Detail B", "children": []}]}, '
    '{"topic": "Subtopic 2", "children": []}]}'
)


class MindMapError(Exception):
    """Raised when the model's output can't be parsed/validated as a mind map."""
    pass


def _extract_json_object_mindmap(text: str) -> dict:
    text = text.strip()
    fence_match = re.search(r"```(?:json)?\s*(.*?)```", text, re.DOTALL)
    if fence_match:
        text = fence_match.group(1).strip()

    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end < start:
        preview = text[:200].replace("\n", " ") if text else "(empty response)"
        raise MindMapError(
            f"The model didn't respond with JSON — it wrote this instead: \"{preview}"
            f"{'...' if len(text) > 200 else ''}\""
        )
    json_str = text[start:end + 1]
    try:
        return json.loads(json_str, strict=False)
    except json.JSONDecodeError:
        pass

    try:
        return json.loads(repair_json(json_str), strict=False)
    except (json.JSONDecodeError, Exception) as e:
        raise MindMapError(
            f"Couldn't parse the model's output as JSON, even after attempting repair: {e}. "
            f"Try again, possibly with a narrower topic."
        )


def _validate_tree(node, depth: int = 0, max_depth: int = 4) -> dict:
    """Structural validation + normalization: every node gets a string 'topic' and
    a list 'children', recursively. Caps runaway depth defensively."""
    if not isinstance(node, dict):
        raise MindMapError("Malformed mind map: expected an object at each node.")
    topic = str(node.get("topic", "")).strip()
    if not topic:
        raise MindMapError("Malformed mind map: a node is missing its 'topic'.")

    raw_children = node.get("children", [])
    if not isinstance(raw_children, list):
        raw_children = []

    children = []
    if depth < max_depth:
        for child in raw_children:
            children.append(_validate_tree(child, depth + 1, max_depth))

    return {"topic": topic, "children": children}


def generate_mindmap(topic: str, source_text: str = None, max_children: int = 5) -> dict:
    """Generate a hierarchical mind map.

    Args:
        topic: the subject to map out (e.g. "Machine Learning", or a question).
        source_text: optional grounding text (e.g. an uploaded document's content).
            If provided, the map is built from this content rather than general
            knowledge.
        max_children: soft cap on branches per node, passed to the model as guidance.

    Returns:
        A nested dict: {"topic": str, "children": [ ...same shape... ]}

    Raises:
        ValueError: empty/missing topic
        LLMError: no provider available or the provider call failed
        MindMapError: the model's output couldn't be parsed/validated
    """
    if not topic or not topic.strip():
        raise ValueError("A topic is required to generate a mind map.")

    if source_text:
        truncated = source_text[:MAX_SOURCE_CHARS_MINDMAP]
        user_prompt = (
            f"Document content:\n\n{truncated}\n\n"
            f'Create a mind map of the key concepts in this document, organized around '
            f'the theme: "{topic}". Use a root topic, up to {max_children} main branches, '
            f"and up to {max_children} sub-branches under each, going up to 3 levels deep. "
            f"Base it strictly on the document content.\n\n"
            f"Respond with ONLY a valid JSON object matching this schema, no prose, no "
            f"code fences:\n{SCHEMA_HINT_MINDMAP}"
        )
    else:
        user_prompt = (
            f'Create a mind map explaining: "{topic}". Use a root topic, up to '
            f"{max_children} main branches, and up to {max_children} sub-branches under "
            f"each, going up to 3 levels deep.\n\n"
            f"Respond with ONLY a valid JSON object matching this schema, no prose, no "
            f"code fences:\n{SCHEMA_HINT_MINDMAP}"
        )

    raw = generate_raw(SYSTEM_PROMPT_MINDMAP, user_prompt, max_tokens=2500)
    try:
        tree = _extract_json_object_mindmap(raw)
    except MindMapError:
        # The model occasionally responds with prose instead of JSON. Retry
        # once with a more insistent reminder before giving up.
        reinforced_prompt = (
            user_prompt + "\n\nREMINDER: your entire response must be ONLY the JSON "
            "object described above — no commentary, no explanation, before or after."
        )
        raw = generate_raw(SYSTEM_PROMPT_MINDMAP, reinforced_prompt, max_tokens=2500)
        tree = _extract_json_object_mindmap(raw)  # let this raise if it fails again
    return _validate_tree(tree)


def tree_to_graph(tree: dict):
    """Flatten a mind-map tree into (nodes, edges) lists suitable for a graph
    visualization library.

    nodes: [{"id": int, "label": str, "level": int}]
    edges: [{"from": int, "to": int}]
    """
    nodes = []
    edges = []
    counter = {"next_id": 0}

    def walk(node, parent_id, level):
        my_id = counter["next_id"]
        counter["next_id"] += 1
        nodes.append({"id": my_id, "label": node["topic"], "level": level})
        if parent_id is not None:
            edges.append({"from": parent_id, "to": my_id})
        for child in node.get("children", []):
            walk(child, my_id, level + 1)

    walk(tree, None, 0)
    return nodes, edges


def render_outline(tree: dict, indent: int = 0) -> str:
    """Render the tree as a markdown bullet-list outline — a plain-text fallback
    view alongside the interactive graph (also handy for copy/paste)."""
    bullet = "  " * indent + f"- {tree['topic']}"
    lines = [bullet]
    for child in tree.get("children", []):
        lines.append(render_outline(child, indent + 1))
    return "\n".join(lines)

# ============================================================================
# --- from doc_comparison.py ===========================================
# ============================================================================

import json
import re

from json_repair import repair_json


MAX_CHARS_PER_DOC = 6000  # keep combined prompt within free-tier context limits

SYSTEM_PROMPT_DOCCOMPARE = (
    "You are a meticulous document comparison assistant. You output ONLY a valid JSON "
    "object, nothing else — no explanations, no markdown formatting, no code fences, "
    "no text before or after the JSON."
)

SCHEMA_HINT_DOCCOMPARE = (
    '{"similarities": ["...", "..."], '
    '"differences": [{"aspect": "...", "document_a": "...", "document_b": "..."}], '
    '"missing_in_a": ["...", "..."], '
    '"missing_in_b": ["...", "..."], '
    '"risks": [{"description": "...", "severity": "Low|Medium|High", "document": "A|B|Both"}]}'
)

VALID_SEVERITIES = {"Low", "Medium", "High"}
VALID_DOC_TAGS = {"A", "B", "Both"}


class DocComparisonError(Exception):
    """Raised when the model's output can't be parsed/validated as a comparison result."""
    pass


def _extract_json_object_doccompare(text: str) -> dict:
    text = text.strip()
    fence_match = re.search(r"```(?:json)?\s*(.*?)```", text, re.DOTALL)
    if fence_match:
        text = fence_match.group(1).strip()

    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end < start:
        preview = text[:200].replace("\n", " ") if text else "(empty response)"
        raise DocComparisonError(
            f"The model didn't respond with JSON — it wrote this instead: \"{preview}"
            f"{'...' if len(text) > 200 else ''}\""
        )
    json_str = text[start:end + 1]
    try:
        return json.loads(json_str, strict=False)
    except json.JSONDecodeError:
        pass

    try:
        return json.loads(repair_json(json_str), strict=False)
    except (json.JSONDecodeError, Exception) as e:
        raise DocComparisonError(
            f"Couldn't parse the model's output as JSON, even after attempting repair: {e}. Try again."
        )


def _as_str_list(value) -> list:
    if not isinstance(value, list):
        return []
    result = []
    for v in value:
        if v is None:
            continue
        s = str(v).strip()
        if s:
            result.append(s)
    return result


def _validate_result(data: dict) -> dict:
    if not isinstance(data, dict):
        raise DocComparisonError("Malformed comparison result: expected a JSON object.")

    similarities = _as_str_list(data.get("similarities"))
    missing_in_a = _as_str_list(data.get("missing_in_a"))
    missing_in_b = _as_str_list(data.get("missing_in_b"))

    differences = []
    for d in data.get("differences", []) if isinstance(data.get("differences"), list) else []:
        if not isinstance(d, dict):
            continue
        aspect = str(d.get("aspect", "")).strip()
        if not aspect:
            continue
        differences.append({
            "aspect": aspect,
            "document_a": str(d.get("document_a", "")).strip(),
            "document_b": str(d.get("document_b", "")).strip(),
        })

    risks = []
    for r in data.get("risks", []) if isinstance(data.get("risks"), list) else []:
        if not isinstance(r, dict):
            continue
        description = str(r.get("description", "")).strip()
        if not description:
            continue
        severity = str(r.get("severity", "Medium")).strip().capitalize()
        if severity not in VALID_SEVERITIES:
            severity = "Medium"
        doc_tag = str(r.get("document", "Both")).strip().capitalize()
        if doc_tag not in VALID_DOC_TAGS:
            doc_tag = "Both"
        risks.append({"description": description, "severity": severity, "document": doc_tag})

    return {
        "similarities": similarities,
        "differences": differences,
        "missing_in_a": missing_in_a,
        "missing_in_b": missing_in_b,
        "risks": risks,
    }


def compare_documents(text_a: str, name_a: str, text_b: str, name_b: str) -> dict:
    """Compare two documents and return a structured breakdown.

    Returns a dict with keys: similarities (list[str]), differences
    (list[{aspect, document_a, document_b}]), missing_in_a (list[str]),
    missing_in_b (list[str]), risks (list[{description, severity, document}]).

    Raises:
        ValueError: either document text is empty
        LLMError: no provider available or the provider call failed
        DocComparisonError: the model's output couldn't be parsed/validated
    """
    if not text_a or not text_a.strip():
        raise ValueError(f"'{name_a}' has no extracted text to compare.")
    if not text_b or not text_b.strip():
        raise ValueError(f"'{name_b}' has no extracted text to compare.")

    truncated_a = text_a[:MAX_CHARS_PER_DOC]
    truncated_b = text_b[:MAX_CHARS_PER_DOC]

    user_prompt = (
        f"Document A ({name_a}):\n{truncated_a}\n\n"
        f"Document B ({name_b}):\n{truncated_b}\n\n"
        "Compare these two documents. Identify:\n"
        "1. Key similarities between them\n"
        "2. Key differences, aspect by aspect (e.g. terms, clauses, findings, values)\n"
        "3. Anything present in Document A but missing from Document B\n"
        "4. Anything present in Document B but missing from Document A\n"
        "5. Any risks worth flagging (e.g. unfavorable terms, contradictions, gaps, "
        "ambiguities), each with a severity (Low/Medium/High) and which document(s) "
        "it applies to (A, B, or Both)\n\n"
        "If the documents are identical or nearly identical, that's fine — just "
        "reflect that in the JSON (e.g. list the shared content under similarities, "
        "with empty differences/missing_in_a/missing_in_b/risks arrays). Do NOT "
        "respond with a plain-text explanation instead of the JSON schema, even in "
        "that case.\n\n"
        f"Respond with ONLY a valid JSON object matching this schema, no prose, no "
        f"code fences:\n{SCHEMA_HINT_DOCCOMPARE}"
    )

    raw = generate_raw(SYSTEM_PROMPT_DOCCOMPARE, user_prompt, max_tokens=3800)
    try:
        data = _extract_json_object_doccompare(raw)
    except DocComparisonError:
        # The model likely wrote prose instead of JSON (common when documents are
        # very similar/identical). Retry once with a more insistent reminder
        # before giving up — cheap compared to making the user click again.
        reinforced_prompt = (
            user_prompt + "\n\nREMINDER: your entire response must be ONLY the JSON "
            "object described above — no commentary, no explanation, before or after."
        )
        raw = generate_raw(SYSTEM_PROMPT_DOCCOMPARE, reinforced_prompt, max_tokens=3800)
        data = _extract_json_object_doccompare(raw)  # let this raise (with the raw preview) if it fails again
    return _validate_result(data)

# ============================================================================
# --- from deep_research.py ============================================
# ============================================================================

import json
import re

from json_repair import repair_json


MAX_SUBQUERIES = 6
DEFAULT_SUBQUERIES = 4
DEFAULT_TOP_K_PER_SUBQUERY = 3

VALID_CONFIDENCE_LABELS = {"Low", "Medium", "High"}

PLAN_SYSTEM_PROMPT = (
    "You are a research planning assistant. You output ONLY a valid JSON array of "
    "strings, nothing else — no explanations, no markdown, no code fences."
)

SYNTHESIS_SYSTEM_PROMPT = (
    "You are a rigorous deep-research assistant. Reason step by step using ONLY the "
    "provided source excerpts before answering. You output ONLY a valid JSON object, "
    "nothing else — no explanations, no markdown, no code fences, no text before or "
    "after the JSON."
)


class DeepResearchError(Exception):
    """Raised when a stage of the pipeline produces output that can't be parsed/validated."""
    pass


# ---------------------------------------------------------------------------
# JSON parsing helpers (duplicated intentionally, mirrors the pattern used in
# quiz_generator.py / mindmap_generator.py / doc_comparison.py — keeps each
# feature module independent so a bug in one can't ripple into the others)
# ---------------------------------------------------------------------------

def _extract_json_array(text: str) -> list:
    text = text.strip()
    fence_match = re.search(r"```(?:json)?\s*(.*?)```", text, re.DOTALL)
    if fence_match:
        text = fence_match.group(1).strip()
    start = text.find("[")
    end = text.rfind("]")
    if start == -1 or end == -1 or end < start:
        preview = text[:200].replace("\n", " ") if text else "(empty response)"
        raise DeepResearchError(
            f"The model didn't respond with JSON — it wrote this instead: \"{preview}"
            f"{'...' if len(text) > 200 else ''}\""
        )
    json_str = text[start:end + 1]
    try:
        return json.loads(json_str, strict=False)
    except json.JSONDecodeError:
        pass

    try:
        return json.loads(repair_json(json_str), strict=False)
    except (json.JSONDecodeError, Exception) as e:
        raise DeepResearchError(f"Couldn't parse the search plan as JSON, even after attempting repair: {e}")


def _extract_json_object_deepresearch(text: str) -> dict:
    text = text.strip()
    fence_match = re.search(r"```(?:json)?\s*(.*?)```", text, re.DOTALL)
    if fence_match:
        text = fence_match.group(1).strip()
    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end < start:
        preview = text[:200].replace("\n", " ") if text else "(empty response)"
        raise DeepResearchError(
            f"The model didn't respond with JSON — it wrote this instead: \"{preview}"
            f"{'...' if len(text) > 200 else ''}\""
        )
    json_str = text[start:end + 1]
    try:
        return json.loads(json_str, strict=False)
    except json.JSONDecodeError:
        pass

    try:
        return json.loads(repair_json(json_str), strict=False)
    except (json.JSONDecodeError, Exception) as e:
        raise DeepResearchError(f"Couldn't parse the synthesis result as JSON, even after attempting repair: {e}")


# ---------------------------------------------------------------------------
# Stage 1: search plan
# ---------------------------------------------------------------------------

def generate_search_plan(question: str, num_subqueries: int = DEFAULT_SUBQUERIES) -> list:
    """Break the question into a handful of focused sub-questions/search queries
    that together would help answer it more thoroughly than a single retrieval pass."""
    if not question or not question.strip():
        raise ValueError("A research question is required.")
    num_subqueries = max(1, min(num_subqueries, MAX_SUBQUERIES))

    user_prompt = (
        f"Break this research question down into exactly {num_subqueries} focused "
        f'sub-questions or search queries that together would help answer it '
        f'thoroughly:\n\n"{question}"\n\n'
        "Respond with ONLY a JSON array of strings, no prose, no code fences."
    )
    raw = generate_raw(PLAN_SYSTEM_PROMPT, user_prompt, max_tokens=600)
    try:
        items = _extract_json_array(raw)
    except DeepResearchError:
        # The model occasionally responds with prose instead of JSON. Retry
        # once with a more insistent reminder before giving up.
        reinforced_prompt = (
            user_prompt + "\n\nREMINDER: your entire response must be ONLY the JSON "
            "array described above — no commentary, no explanation, before or after."
        )
        raw = generate_raw(PLAN_SYSTEM_PROMPT, reinforced_prompt, max_tokens=600)
        items = _extract_json_array(raw)  # let this raise if it fails again
    plan = [str(x).strip() for x in items if isinstance(x, (str, int, float)) and str(x).strip()]
    if not plan:
        raise DeepResearchError("The model returned an empty search plan.")
    return plan


# ---------------------------------------------------------------------------
# Stage 2: multi-query retrieval + dedup (pure logic, no LLM call — testable
# without mocking anything beyond a fake vector store)
# ---------------------------------------------------------------------------

def retrieve_for_plan(vector_store, plan: list, top_k_per_subquery: int = DEFAULT_TOP_K_PER_SUBQUERY) -> list:
    """Search the vector store once per sub-question in the plan, and merge results
    into a deduplicated, numbered source list. The same chunk surfaced by multiple
    sub-questions is recorded once, with all the sub-questions that found it."""
    seen = {}
    order = []
    for subq in plan:
        results = vector_store.search(subq, top_k=top_k_per_subquery)
        for r in results:
            dedup_key = (r.filename, r.page, r.text)
            if dedup_key not in seen:
                seen[dedup_key] = {
                    "filename": r.filename, "page": r.page, "text": r.text,
                    "score": r.score, "subqueries": [subq],
                }
                order.append(dedup_key)
            elif subq not in seen[dedup_key]["subqueries"]:
                seen[dedup_key]["subqueries"].append(subq)

    sources = [seen[k] for k in order]
    for i, s in enumerate(sources, start=1):
        s["index"] = i
    return sources


def _build_numbered_context(sources: list) -> str:
    parts = []
    for s in sources:
        parts.append(f"[Source {s['index']}] ({s['filename']}, page {s['page']}):\n{s['text']}")
    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Stage 3: step-by-step synthesis with confidence scoring
# ---------------------------------------------------------------------------

def _validate_confidence(raw_confidence) -> dict:
    if not isinstance(raw_confidence, dict):
        raw_confidence = {}
    try:
        score = int(round(float(raw_confidence.get("score", 50))))
    except (TypeError, ValueError):
        score = 50
    score = max(0, min(100, score))

    label = str(raw_confidence.get("label", "")).strip().capitalize()
    if label not in VALID_CONFIDENCE_LABELS:
        # derive a sensible label from the score if the model didn't give a valid one
        label = "High" if score >= 75 else "Medium" if score >= 40 else "Low"

    justification = str(raw_confidence.get("justification", "")).strip()
    return {"score": score, "label": label, "justification": justification}


def _synthesize(question: str, plan: list, sources: list) -> dict:
    context = _build_numbered_context(sources)
    plan_text = "\n".join(f"- {p}" for p in plan)
    user_prompt = (
        f"Original question: {question}\n\n"
        f"Research plan (sub-questions investigated):\n{plan_text}\n\n"
        f"Source excerpts:\n\n{context}\n\n"
        "Reason through this step by step across the sub-questions, then provide a "
        "final answer. Cite sources inline as [Source N]. Assess your confidence "
        "honestly — lower it if sources conflict, are incomplete, or are only "
        "tangentially relevant to the question.\n\n"
        "Respond with ONLY a valid JSON object matching this schema, no prose, no "
        "code fences:\n"
        '{"reasoning_steps": ["step 1 ...", "step 2 ..."], '
        '"answer": "... with [Source N] citations ...", '
        '"confidence": {"score": 0-100, "label": "Low|Medium|High", "justification": "..."}}'
    )
    raw = generate_raw(SYNTHESIS_SYSTEM_PROMPT, user_prompt, max_tokens=3200)
    try:
        data = _extract_json_object_deepresearch(raw)
    except DeepResearchError:
        # The model occasionally responds with prose instead of JSON. Retry
        # once with a more insistent reminder before giving up.
        reinforced_prompt = (
            user_prompt + "\n\nREMINDER: your entire response must be ONLY the JSON "
            "object described above — no commentary, no explanation, before or after."
        )
        raw = generate_raw(SYNTHESIS_SYSTEM_PROMPT, reinforced_prompt, max_tokens=3200)
        data = _extract_json_object_deepresearch(raw)  # let this raise if it fails again

    reasoning_steps = data.get("reasoning_steps", [])
    if not isinstance(reasoning_steps, list):
        reasoning_steps = []
    reasoning_steps = [str(s).strip() for s in reasoning_steps if str(s).strip()]

    answer = str(data.get("answer", "")).strip()
    if not answer:
        raise DeepResearchError("The model didn't return a final answer.")

    confidence = _validate_confidence(data.get("confidence"))

    return {"reasoning_steps": reasoning_steps, "answer": answer, "confidence": confidence}


# ---------------------------------------------------------------------------
# Full pipeline
# ---------------------------------------------------------------------------

def run_deep_research(
    question: str,
    vector_store,
    num_subqueries: int = DEFAULT_SUBQUERIES,
    top_k_per_subquery: int = DEFAULT_TOP_K_PER_SUBQUERY,
) -> dict:
    """Run the full deep research pipeline and return a structured result:

        {
            "search_plan": [str, ...],
            "reasoning_steps": [str, ...],
            "answer": str,
            "confidence": {"score": int, "label": str, "justification": str},
            "sources": [{"index": int, "filename": str, "page": int, "text": str,
                         "score": float, "subqueries": [str, ...]}, ...],
        }

    Raises:
        ValueError: empty question, or no documents indexed yet
        LLMError: no provider available or a provider call failed
        DeepResearchError: a pipeline stage's output couldn't be parsed/validated
    """
    if not question or not question.strip():
        raise ValueError("A research question is required.")
    if not vector_store.has_documents():
        raise ValueError("No documents indexed yet — upload at least one PDF first.")

    plan = generate_search_plan(question, num_subqueries)
    sources = retrieve_for_plan(vector_store, plan, top_k_per_subquery)

    if not sources:
        return {
            "search_plan": plan,
            "reasoning_steps": [],
            "answer": "No relevant content was found in the uploaded documents for this question.",
            "confidence": {
                "score": 0, "label": "Low",
                "justification": "No supporting evidence was retrieved from the uploaded documents.",
            },
            "sources": [],
        }

    synthesis = _synthesize(question, plan, sources)
    synthesis["search_plan"] = plan
    synthesis["sources"] = sources
    return synthesis

# ============================================================================
# --- from translator.py ===============================================
# ============================================================================

MAX_CHARS = 8000  # keep translation requests within free-tier context limits

SYSTEM_PROMPT_TRANSLATOR = (
    "You are a precise, fluent translator. Translate the user's text into the "
    "requested target language. Preserve meaning, tone, and formatting (line "
    "breaks, lists) as closely as possible. Output ONLY the translated text — "
    "no commentary, no explanations, no notes about the translation."
)


def translate_text(text: str, target_language: str) -> str:
    """Translate text into target_language.

    Raises:
        ValueError: empty text or empty target language
        LLMError: no provider available or the provider call failed
    """
    if not text or not text.strip():
        raise ValueError("There's no text to translate.")
    if not target_language or not target_language.strip():
        raise ValueError("A target language is required.")

    truncated = text[:MAX_CHARS]
    truncated_note = ""
    if len(text) > MAX_CHARS:
        truncated_note = f"\n\n*(Note: input was truncated to the first {MAX_CHARS} characters.)*"

    user_prompt = f"Translate the following text into {target_language.strip()}:\n\n{truncated}"
    result = generate_raw(SYSTEM_PROMPT_TRANSLATOR, user_prompt, max_tokens=3000)
    return result + truncated_note

# ============================================================================
# --- from voice_client.py =============================================
# ============================================================================

import io
import os
import re
import wave


STT_API_URL = "https://api.groq.com/openai/v1/audio/transcriptions"
TTS_API_URL = "https://api.groq.com/openai/v1/audio/speech"

STT_MODEL = "whisper-large-v3-turbo"
TTS_MODEL = "canopylabs/orpheus-v1-english"

TTS_MAX_CHARS_PER_CALL = 200  # hard limit imposed by the Orpheus API
TTS_DEFAULT_MAX_TOTAL_CHARS = 600  # cap total chars spoken per "read aloud" (cost/rate-limit control)
TTS_DEFAULT_VOICE = "troy"
TTS_VALID_VOICES = {"autumn", "diana", "hannah", "austin", "daniel", "troy"}


def transcribe_audio(audio_bytes: bytes, filename: str = "audio.wav") -> str:
    """Transcribe recorded speech to text using Groq's Whisper API.

    Raises:
        LLMError: no Groq API key configured, or the transcription call failed
    """
    groq_key = os.environ.get("GROQ_API_KEY")
    if not groq_key:
        raise LLMError(
            "Voice input requires a Groq API key (GROQ_API_KEY) — Groq is currently "
            "the only configured provider with speech-to-text support in this app."
        )
    files = {"file": (filename, audio_bytes)}
    data = {"model": STT_MODEL, "response_format": "json"}
    resp = _groq_post_multipart(STT_API_URL, groq_key, files=files, data=data, timeout=60)
    if resp.status_code != 200:
        raise LLMError(f"Groq transcription error {resp.status_code}: {resp.text[:300]}")
    text = resp.json().get("text", "").strip()
    if not text:
        raise LLMError("No speech was detected in the recording. Try again and speak clearly.")
    return text


def _split_for_tts(text: str, max_chars: int = TTS_MAX_CHARS_PER_CALL) -> list:
    """Split text into chunks that each fit within Orpheus's per-request character
    limit, breaking on sentence boundaries where possible so each chunk sounds
    natural rather than cutting off mid-sentence."""
    text = text.strip()
    if not text:
        return []

    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks = []
    current = ""

    for sentence in sentences:
        if len(sentence) > max_chars:
            if current:
                chunks.append(current)
                current = ""
            for i in range(0, len(sentence), max_chars):
                chunks.append(sentence[i:i + max_chars])
            continue

        candidate = f"{current} {sentence}".strip() if current else sentence
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                chunks.append(current)
            current = sentence

    if current:
        chunks.append(current)
    return chunks


def _concatenate_wavs(wav_byte_chunks: list) -> bytes:
    """Concatenate multiple WAV clips (same format) into a single playable WAV."""
    if len(wav_byte_chunks) == 1:
        return wav_byte_chunks[0]

    with wave.open(io.BytesIO(wav_byte_chunks[0]), "rb") as first:
        params = first.getparams()
        frames = [first.readframes(first.getnframes())]

    for wav_bytes in wav_byte_chunks[1:]:
        with wave.open(io.BytesIO(wav_bytes), "rb") as w:
            frames.append(w.readframes(w.getnframes()))

    out_buf = io.BytesIO()
    with wave.open(out_buf, "wb") as out:
        out.setparams(params)
        for f in frames:
            out.writeframes(f)
    return out_buf.getvalue()


def synthesize_speech(
    text: str,
    voice: str = TTS_DEFAULT_VOICE,
    max_total_chars: int = TTS_DEFAULT_MAX_TOTAL_CHARS,
) -> bytes:
    """Convert text to speech using Groq's Orpheus TTS model. Text longer than
    max_total_chars is truncated first (to control cost and API call count —
    see module docstring re: TTS not being confirmed free-tier), then split
    into <=200-char chunks (Orpheus's hard limit) and the resulting audio
    clips are concatenated into one WAV file.

    Raises:
        ValueError: empty text
        LLMError: no Groq API key configured, or a synthesis call failed
    """
    if not text or not text.strip():
        raise ValueError("There's no text to speak.")

    groq_key = os.environ.get("GROQ_API_KEY")
    if not groq_key:
        raise LLMError(
            "Voice output requires a Groq API key (GROQ_API_KEY) — Groq is currently "
            "the only configured provider with text-to-speech support in this app."
        )

    if voice not in TTS_VALID_VOICES:
        voice = TTS_DEFAULT_VOICE

    truncated = text.strip()[:max_total_chars]
    chunks = _split_for_tts(truncated)
    if not chunks:
        raise ValueError("There's no text to speak.")

    wav_parts = []
    for chunk in chunks:
        payload = {"model": TTS_MODEL, "input": chunk, "voice": voice, "response_format": "wav"}
        resp = _groq_post(payload, groq_key, url=TTS_API_URL, timeout=60)
        if resp.status_code != 200:
            raise LLMError(f"Groq TTS error {resp.status_code}: {resp.text[:300]}")
        wav_parts.append(resp.content)

    return _concatenate_wavs(wav_parts)

# ============================================================================
# --- from export_utils.py =============================================
# ============================================================================

from datetime import datetime
from pathlib import Path

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable

EXPORT_DIR = Path(__file__).parent / "exports"


def _escape(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def export_chat_to_pdf(chat_history: list, session_id: str) -> str:
    """Render the given chat history into a PDF and return the file path."""
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    filename = f"research_qa_{session_id}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pdf"
    filepath = EXPORT_DIR / filename

    styles = getSampleStyleSheet()
    title_style = styles["Title"]
    question_style = ParagraphStyle(
        "Question", parent=styles["Heading3"], textColor=colors.HexColor("#1a4d8f"), spaceBefore=14
    )
    answer_style = ParagraphStyle("Answer", parent=styles["Normal"], spaceBefore=6, leading=15)
    source_style = ParagraphStyle(
        "Source", parent=styles["Normal"], fontSize=8.5, textColor=colors.HexColor("#555555"),
        leftIndent=14, spaceBefore=2,
    )
    meta_style = ParagraphStyle("Meta", parent=styles["Normal"], fontSize=9, textColor=colors.grey)

    doc = SimpleDocTemplate(str(filepath), pagesize=letter,
                             topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    story = [
        Paragraph("AI Research Assistant — Q&A Export", title_style),
        Paragraph(f"Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}", meta_style),
        Spacer(1, 16),
    ]

    for i, entry in enumerate(chat_history, start=1):
        story.append(Paragraph(f"Q{i}: {_escape(entry['question'])}", question_style))
        story.append(Paragraph(_escape(entry["answer"]).replace("\n", "<br/>"), answer_style))

        sources = entry.get("sources") or []
        if sources:
            story.append(Spacer(1, 4))
            for s in sources:
                label = f"Source: {_escape(s.get('filename', ''))}, page {s.get('page', '?')}"
                story.append(Paragraph(label, source_style))

        story.append(Spacer(1, 8))
        story.append(HRFlowable(width="100%", color=colors.HexColor("#dddddd")))

    doc.build(story)
    return str(filepath)


# ---------------------------------------------------------------------------
# Expose database.py's functions as db.<name>(...) so every db.xxx() call in
# the app section below (originally `import database as db`) keeps working
# unchanged, without needing a separate module.
# ---------------------------------------------------------------------------
db = types.SimpleNamespace(
    init_db=init_db,
    add_document=add_document,
    get_documents=get_documents,
    add_chat_entry=add_chat_entry,
    get_chat_history=get_chat_history,
    clear_chat_history=clear_chat_history,
    clear_session=clear_session,
    log_tool_usage=log_tool_usage,
    get_tool_usage_counts=get_tool_usage_counts,
)


# ============================================================================
# --- from app.py (main Streamlit script) ==============================
# ============================================================================

ENV_PATH = Path(__file__).parent / ".env"
if not ENV_PATH.exists():
    ENV_PATH.touch()
load_dotenv(ENV_PATH, override=True)

st.set_page_config(page_title="AI Research Assistant", page_icon="📚", layout="wide")

# Persistent storage for uploaded PDFs, keyed by session, so features that need
# the original file (e.g. rendering a page as an image for chart analysis) can
# re-open it later — the RAG index only stores extracted text, not the file itself.
UPLOADS_DIR = Path(__file__).parent / "data" / "uploads"


def _uploaded_file_path(session_id: str, filename: str) -> Path:
    return UPLOADS_DIR / session_id / filename


# ---------------------------------------------------------------------------
# AI Tools registry — each entry drives both the sidebar nav and the main-area
# header. Only tools that are actually implemented are listed here.
# ---------------------------------------------------------------------------
TOOLS = [
    {"key": "dashboard", "label": "📊 Dashboard",
     "caption": "An overview of this session's activity across every tool."},
    {"key": "chat", "label": "💬 Chat with PDF",
     "caption": "Semantic search across all uploaded PDFs, with cited, highlighted sources."},
    {"key": "image", "label": "🖼️ Chat with Image",
     "caption": "Ask a vision model to read charts, diagrams, tables, and photos on any page."},
    {"key": "quiz", "label": "🎯 Quiz & Flashcards",
     "caption": "Generate MCQs, coding questions, interview questions, or flashcards."},
    {"key": "mindmap", "label": "🧠 Mind Map",
     "caption": "Visualize any topic — or a document's key concepts — as an interactive graph."},
    {"key": "compare", "label": "⚖️ Compare Documents",
     "caption": "Find similarities, differences, missing clauses, and risks between two documents."},
    {"key": "research", "label": "🔬 Deep Research",
     "caption": "Search plan, step-by-step reasoning, and an honest confidence score."},
    {"key": "translate", "label": "🌍 Translate",
     "caption": "Translate pasted text or an uploaded document into any language."},
    {"key": "voice", "label": "🎤 Voice Chat",
     "caption": "Ask your documents a question out loud instead of typing."},
]

# ---------------------------------------------------------------------------
# Session state / persistence setup
# ---------------------------------------------------------------------------
db.init_db()

if "session_id" not in st.session_state:
    st.session_state.session_id = str(uuid.uuid4())[:8]

if "session_start" not in st.session_state:
    st.session_state.session_start = datetime.utcnow()

if "vector_store" not in st.session_state:
    st.session_state.vector_store = VectorStore(st.session_state.session_id)

if "chat_history" not in st.session_state:
    st.session_state.chat_history = db.get_chat_history(st.session_state.session_id)

vs: VectorStore = st.session_state.vector_store


def _render_mindmap_html(tree: dict) -> str:
    """Build a self-contained HTML snippet that renders the given mind-map tree as
    an interactive, draggable graph using vis-network (loaded from a CDN)."""
    nodes, edges = tree_to_graph(tree)

    level_colors = {
        0: {"background": "#d97706", "border": "#92400e"},  # root — amber
        1: {"background": "#2563eb", "border": "#1e3a8a"},  # branches — blue
    }
    default_color = {"background": "#16a34a", "border": "#14532d"}  # leaves — green

    vis_nodes = []
    for n in nodes:
        color = level_colors.get(n["level"], default_color)
        vis_nodes.append({
            "id": n["id"],
            "label": n["label"],
            "level": n["level"],
            "color": color,
            "font": {"color": "#ffffff", "size": 15 if n["level"] == 0 else 13},
        })

    nodes_json = json.dumps(vis_nodes).replace("</", "<\\/")
    edges_json = json.dumps(edges).replace("</", "<\\/")

    return f"""
    <div id="mindmap" style="height: 520px; border: 1px solid #444; border-radius: 10px; background: #0e1117;"></div>
    <script src="https://unpkg.com/vis-network@9.1.6/standalone/umd/vis-network.min.js"></script>
    <script>
      const nodes = new vis.DataSet({nodes_json});
      const edges = new vis.DataSet({edges_json});
      const container = document.getElementById('mindmap');
      const data = {{ nodes: nodes, edges: edges }};
      const options = {{
        layout: {{
          hierarchical: {{
            direction: 'UD',
            sortMethod: 'directed',
            nodeSpacing: 160,
            levelSeparation: 130
          }}
        }},
        nodes: {{
          shape: 'box',
          margin: 10,
          borderWidth: 2,
          shapeProperties: {{ borderRadius: 6 }}
        }},
        edges: {{
          color: {{ color: '#555555', highlight: '#999999' }},
          smooth: {{ type: 'cubicBezier', forceDirection: 'vertical' }},
          width: 1.5
        }},
        physics: false,
        interaction: {{ dragNodes: true, zoomView: true, dragView: true }}
      }};
      new vis.Network(container, data, options);
    </script>
    """


def _render_pdf_viewer_html(pdf_bytes: bytes, height: int = 600) -> str:
    """Embed a PDF directly via a base64 data URI so the whole document can be
    viewed inline. Deliberately not using Streamlit's built-in st.pdf — that
    command requires installing a separate `streamlit-pdf` package as an extra
    even on Streamlit versions where it exists, which is exactly the kind of
    hidden dependency that's caused problems before. This approach needs
    nothing beyond what's already in requirements.txt, and works in any
    browser with native PDF viewing support (all modern browsers)."""
    b64 = base64.b64encode(pdf_bytes).decode("utf-8")
    return (
        f'<embed src="data:application/pdf;base64,{b64}" '
        f'width="100%" height="{height}" type="application/pdf" />'
    )


@st.fragment(run_every="1s")
def _render_live_session_timer(start_time: datetime):
    """Renders the session duration as its own fragment that reruns every
    second independently of the rest of the page — this is what makes it
    tick continuously like a real clock instead of only updating whenever
    something else on the page happens to trigger a full rerun."""
    elapsed = datetime.utcnow() - start_time
    total_seconds = int(elapsed.total_seconds())
    hours, remainder = divmod(total_seconds, 3600)
    minutes, seconds = divmod(remainder, 60)
    duration_str = f"{hours}h {minutes}m {seconds}s" if hours > 0 else f"{minutes}m {seconds}s"
    st.metric("Session Time", duration_str)


# ---------------------------------------------------------------------------
# Sidebar — uploads, document list, AI Tools nav, settings
# ---------------------------------------------------------------------------
with st.sidebar:
    st.title("📚 Research Assistant")
    st.caption(f"Session: `{st.session_state.session_id}`")

    st.markdown("### 📂 Upload Knowledge Sources")
    with st.container(border=True):
        st.markdown("**Supported files**")
        st.markdown(
            "📄 PDF · 📃 DOCX · 📽️ PPTX · 📊 XLSX · 📈 CSV  \n"
            "📝 TXT / Markdown · 💻 Code files  \n"
            "🖼️ Images (JPG, PNG, WEBP, BMP, TIFF)"
        )
        uploader_types = sorted(ext.lstrip(".") for ext in SUPPORTED_EXTENSIONS)
        uploaded_files = st.file_uploader(
            "Drag & drop files here, or click to browse",
            type=uploader_types, accept_multiple_files=True,
        )
        st.caption("Max 200MB per file · images are described via a vision model (needs Groq)")

    if uploaded_files:
        for uf in uploaded_files:
            already_indexed = uf.name in vs.list_filenames()
            if already_indexed:
                continue
            with st.spinner(f"Indexing {uf.name}..."):
                persistent_path = _uploaded_file_path(st.session_state.session_id, uf.name)
                persistent_path.parent.mkdir(parents=True, exist_ok=True)
                persistent_path.write_bytes(uf.getvalue())
                try:
                    num_chunks, num_pages = vs.add_document(str(persistent_path), uf.name)
                    if num_chunks == 0:
                        st.warning(f"{uf.name}: no extractable text found — skipped.")
                        persistent_path.unlink(missing_ok=True)
                    else:
                        db.add_document(uf.name, num_pages, num_chunks, st.session_state.session_id)
                        st.success(f"{uf.name}: {num_pages} page(s), {num_chunks} chunks indexed")
                except (DocumentLoadError, LLMError) as e:
                    st.error(f"{uf.name}: {e}")
                    persistent_path.unlink(missing_ok=True)
                except Exception:
                    persistent_path.unlink(missing_ok=True)
                    raise

    docs = db.get_documents(st.session_state.session_id)
    st.markdown(f"### 📁 Indexed Documents ({len(docs)})")
    if docs:
        with st.container(border=True):
            for d in docs:
                st.markdown(f"📄 **{d['filename']}**")
                st.caption(f"{d['num_pages']} pages · {d['num_chunks']} chunks")
    else:
        st.caption("No documents yet. Upload a PDF above to get started.")

    st.divider()
    st.subheader("🧰 AI Tools")
    tool_labels = [t["label"] for t in TOOLS]
    selected_label = st.radio(
        "Choose a tool", tool_labels, key="active_tool_label", label_visibility="collapsed"
    )
    active_tool = next(t["key"] for t in TOOLS if t["label"] == selected_label)

    st.divider()

    groq_ready = bool(os.environ.get("GROQ_API_KEY"))
    hf_ready = bool(os.environ.get("HF_API_TOKEN"))

    st.subheader("LLM provider")
    status_lines = [
        f"{'✅' if groq_ready else '⬜'} Groq (`GROQ_API_KEY`)",
        f"{'✅' if hf_ready else '⬜'} Hugging Face (`HF_API_TOKEN`)",
    ]
    st.caption("  \n".join(status_lines))

    if not (groq_ready or hf_ready):
        st.warning(
            "⚠️ No LLM provider available. Set `GROQ_API_KEY` (free, fastest — "
            "console.groq.com) or `HF_API_TOKEN` below to get started.",
            icon="⚠️",
        )

    with st.expander("🔑 Save API key" if not (groq_ready or hf_ready) else "🔑 Update API key"):
        st.caption(
            "Saved to a local `.env` file in the project folder — loads automatically "
            "every time you run the app, no need to set it in the terminal again."
        )
        groq_input = st.text_input(
            "Groq API key", type="password", placeholder="gsk_...",
            help="Get a free key at console.groq.com/keys",
        )
        hf_input = st.text_input(
            "Hugging Face token (optional)", type="password", placeholder="hf_...",
            help="Get a free token at huggingface.co/settings/tokens",
        )
        if st.button("💾 Save", use_container_width=True):
            saved_any = False
            if groq_input.strip():
                set_key(str(ENV_PATH), "GROQ_API_KEY", groq_input.strip())
                os.environ["GROQ_API_KEY"] = groq_input.strip()
                saved_any = True
            if hf_input.strip():
                set_key(str(ENV_PATH), "HF_API_TOKEN", hf_input.strip())
                os.environ["HF_API_TOKEN"] = hf_input.strip()
                saved_any = True
            if saved_any:
                st.success("Saved! This key will now load automatically every time.")
                st.rerun()
            else:
                st.warning("Enter at least one key before saving.")

    top_k = st.slider("Chunks to retrieve per question", 3, 10, 5)

    st.divider()
    col_a, col_b = st.columns(2)
    with col_a:
        if st.button("🗑️ Clear session", use_container_width=True):
            vs.clear()
            db.clear_session(st.session_state.session_id)
            st.session_state.chat_history = []
            session_upload_dir = UPLOADS_DIR / st.session_state.session_id
            if session_upload_dir.exists():
                import shutil
                shutil.rmtree(session_upload_dir, ignore_errors=True)
            st.rerun()
    with col_b:
        if st.button("📄 Export PDF", use_container_width=True, disabled=not st.session_state.chat_history):
            path = export_chat_to_pdf(st.session_state.chat_history, st.session_state.session_id)
            with open(path, "rb") as f:
                st.download_button("⬇️ Download", f, file_name=Path(path).name, mime="application/pdf")

# ---------------------------------------------------------------------------
# Main area — renders exactly one tool, based on the sidebar selection
# ---------------------------------------------------------------------------
active_meta = next(t for t in TOOLS if t["key"] == active_tool)
st.title(active_meta["label"])
st.caption(active_meta["caption"])
st.divider()

# ============================== Dashboard ==============================
if active_tool == "dashboard":
    tool_counts = db.get_tool_usage_counts(st.session_state.session_id)

    total_pages = sum(d["num_pages"] or 0 for d in docs)

    with st.container(border=True):
        st.markdown("##### 📚 Content Overview")
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Documents", len(docs))
        c2.metric("Pages Indexed", total_pages)
        c3.metric("Chunks Indexed", len(vs.chunks))
        with c4:
            _render_live_session_timer(st.session_state.session_start)

    with st.container(border=True):
        st.markdown("##### 💬 Engagement")
        c1, c2, c3 = st.columns(3)
        c1.metric("Questions Asked", len(st.session_state.chat_history))
        c2.metric("Voice Questions", tool_counts.get("voice", 0))
        c3.metric("Translations Run", tool_counts.get("translate", 0))

    with st.container(border=True):
        st.markdown("##### 🛠️ Tools Used")
        c1, c2, c3, c4, c5 = st.columns(5)
        c1.metric("🎯 Quizzes", tool_counts.get("quiz", 0))
        c2.metric("🧠 Mind Maps", tool_counts.get("mindmap", 0))
        c3.metric("⚖️ Comparisons", tool_counts.get("compare", 0))
        c4.metric("🔬 Deep Research", tool_counts.get("research", 0))
        c5.metric("🖼️ Images", tool_counts.get("image", 0))

    st.divider()
    if not docs:
        st.info("Upload a document from the sidebar to get started, then pick a tool on the left.")
    else:
        st.caption("Documents in this session:")
        for d in docs:
            st.markdown(f"- **{d['filename']}** — {d['num_pages']} page(s), {d['num_chunks']} chunks")

# ============================== Chat with PDF ==============================
elif active_tool == "chat":
    pdf_docs = [d for d in docs if d["filename"].lower().endswith(".pdf")]
    if pdf_docs:
        with st.expander("📄 Preview a PDF"):
            preview_doc = st.selectbox(
                "Choose a PDF to preview", [d["filename"] for d in pdf_docs], key="chat_pdf_preview_select",
            )
            preview_path = _uploaded_file_path(st.session_state.session_id, preview_doc)
            if preview_path.exists():
                st.iframe(src=_render_pdf_viewer_html(preview_path.read_bytes()), height=600)
            else:
                st.warning(
                    "The original file for this document isn't available anymore (it was "
                    "uploaded before this feature was added, or the session data was moved). "
                    "Re-upload the PDF to preview it."
                )

    if st.session_state.chat_history:
        if st.button("🗑️ Clear chat"):
            st.session_state.chat_history = []
            db.clear_chat_history(st.session_state.session_id)
            st.rerun()

    for entry in st.session_state.chat_history:
        with st.chat_message("user"):
            st.markdown(entry["question"])
        with st.chat_message("assistant"):
            st.markdown(entry["answer"])
            if entry.get("sources"):
                with st.expander(f"📎 {len(entry['sources'])} source excerpt(s)"):
                    for s in entry["sources"]:
                        st.markdown(f"**{s['filename']}**, page {s['page']} — score {s['score']:.2f}")
                        st.markdown(f"> {s['text']}")

    question = st.chat_input("Ask a question about your uploaded papers...")

    if question:
        if not vs.has_documents():
            st.warning("Please upload at least one PDF first.")
        else:
            with st.chat_message("user"):
                st.markdown(question)

            with st.chat_message("assistant"):
                with st.spinner("Searching documents..."):
                    retrieved = vs.search(question, top_k=top_k)

                try:
                    answer = st.write_stream(generate_answer_stream(question, retrieved))
                except LLMError as e:
                    answer = f"⚠️ {e}"
                    st.markdown(answer)

                sources = [
                    {"filename": r.filename, "page": r.page, "score": r.score, "text": r.text}
                    for r in retrieved
                ]
                if sources:
                    with st.expander(f"📎 {len(sources)} source excerpt(s)"):
                        for s in sources:
                            st.markdown(f"**{s['filename']}**, page {s['page']} — score {s['score']:.2f}")
                            st.markdown(f"> {s['text']}")

            db.add_chat_entry(st.session_state.session_id, question, answer, sources)
            st.session_state.chat_history.append(
                {"question": question, "answer": answer, "sources": sources}
            )

# ============================== Quiz & Flashcards ==============================
elif active_tool == "quiz":
    if not docs:
        st.info("Upload a PDF first to use this tool.")
    else:
        quiz_doc = st.selectbox("Choose a document", [d["filename"] for d in docs], key="quiz_doc_select")
        col1, col2, col3 = st.columns(3)
        with col1:
            quiz_type = st.selectbox("Type", QUIZ_TYPES, key="quiz_type_select")
        with col2:
            quiz_difficulty = st.selectbox("Difficulty", DIFFICULTIES, index=1, key="quiz_difficulty_select")
        with col3:
            quiz_num = st.slider("Items per click", 3, 30, 10, key="quiz_num_select")

        st.caption(
            "Groq's free tier limits how much a single request can generate. Want 50+? "
            "Click Generate a few times with the same type/document — each click adds "
            "more instead of starting over."
        )

        col_gen, col_clear = st.columns([3, 1])
        with col_gen:
            generate_clicked = st.button("➕ Generate", use_container_width=True)
        with col_clear:
            if st.button("🗑️ Clear", use_container_width=True):
                st.session_state.pop("quiz_data", None)
                st.session_state.pop("quiz_checked", None)
                st.rerun()

        if generate_clicked:
            with st.spinner(f"Generating {quiz_num} {quiz_type.lower()}..."):
                matching_chunks = [c for c in vs.chunks if c.filename == quiz_doc]
                full_text = " ".join(c.text for c in matching_chunks)
                try:
                    existing = st.session_state.get("quiz_data")
                    same_batch = existing and existing["type"] == quiz_type and existing["doc"] == quiz_doc
                    offset = existing.get("next_offset", 0) if same_batch else 0

                    new_items = generate_quiz(full_text, quiz_type, quiz_difficulty, quiz_num, offset=offset)

                    # Advance the offset for next click so it explores a different
                    # part of long documents instead of regenerating from the same slice.
                    next_offset = offset + 6000

                    if same_batch:
                        existing["items"].extend(new_items)
                        existing["next_offset"] = next_offset
                    else:
                        st.session_state.quiz_data = {
                            "type": quiz_type, "items": new_items, "doc": quiz_doc,
                            "next_offset": next_offset,
                        }
                    db.log_tool_usage(st.session_state.session_id, "quiz")
                    st.session_state.pop("quiz_checked", None)
                except (LLMError, QuizGenerationError, ValueError) as e:
                    st.error(str(e))

        if st.session_state.get("quiz_data"):
            qdata = st.session_state.quiz_data
            qtype = qdata["type"]
            items = qdata["items"]
            st.markdown(f"**{len(items)} {qtype} — {qdata['doc']}**")

            if qtype == "MCQ":
                for i, item in enumerate(items):
                    st.markdown(f"**Q{i+1}. {item.get('question', '')}**")
                    st.radio(
                        f"mcq_{i}", item.get("options", []), key=f"quiz_mcq_choice_{i}",
                        label_visibility="collapsed",
                    )

                if st.button("✅ Check answers"):
                    st.session_state.quiz_checked = True

                if st.session_state.get("quiz_checked"):
                    score = 0
                    for i, item in enumerate(items):
                        options = item.get("options", [])
                        correct_idx = item.get("correct_index", -1)
                        correct_answer = options[correct_idx] if 0 <= correct_idx < len(options) else None
                        picked = st.session_state.get(f"quiz_mcq_choice_{i}")
                        is_correct = picked == correct_answer
                        score += int(is_correct)
                        icon = "✅" if is_correct else "❌"
                        st.markdown(f"{icon} **Q{i+1}:** correct answer — {correct_answer}")
                        if item.get("explanation"):
                            st.caption(item["explanation"])
                    st.info(f"Score: {score} / {len(items)}")

            elif qtype == "Flashcards":
                for item in items:
                    with st.expander(f"🃏 {item.get('front', '')}"):
                        st.markdown(item.get("back", ""))

            elif qtype == "Coding Questions":
                for i, item in enumerate(items):
                    with st.expander(f"💻 Q{i+1}. {item.get('question', '')}"):
                        if item.get("hint"):
                            st.markdown(f"**Hint:** {item['hint']}")
                        if item.get("sample_approach"):
                            st.markdown(f"**Sample approach:** {item['sample_approach']}")

            elif qtype == "Interview Questions":
                for i, item in enumerate(items):
                    with st.expander(f"🎤 Q{i+1}. {item.get('question', '')}"):
                        if item.get("what_to_cover"):
                            st.markdown(f"**What a strong answer covers:** {item['what_to_cover']}")

# ============================== Mind Map ==============================
elif active_tool == "mindmap":
    mm_topic = st.text_input(
        "Topic or question", placeholder="e.g. Explain Machine Learning",
        key="mindmap_topic_input",
    )
    doc_options = ["(none — use general knowledge)"] + [d["filename"] for d in docs]
    mm_doc_choice = st.selectbox(
        "Base it on a document (optional)", doc_options, key="mindmap_doc_select",
        help="Leave as 'none' to map out any topic from general knowledge, or pick a "
             "document to map its actual content.",
    )

    if st.button("Generate Mind Map"):
        with st.spinner("Building the mind map..."):
            source_text = None
            if mm_doc_choice != doc_options[0]:
                matching_chunks = [c for c in vs.chunks if c.filename == mm_doc_choice]
                source_text = " ".join(c.text for c in matching_chunks)
            try:
                tree = generate_mindmap(mm_topic, source_text=source_text)
                st.session_state.mindmap_data = tree
                db.log_tool_usage(st.session_state.session_id, "mindmap")
            except (ValueError, LLMError, MindMapError) as e:
                st.error(str(e))
                st.session_state.pop("mindmap_data", None)

    if st.session_state.get("mindmap_data"):
        if st.button("🗑️ Clear mind map"):
            st.session_state.pop("mindmap_data", None)
            st.rerun()
        tree = st.session_state.mindmap_data
        st.iframe(src=_render_mindmap_html(tree), height=540)
        with st.expander("📋 View as outline"):
            st.markdown(render_outline(tree))

# ============================== Compare Documents ==============================
elif active_tool == "compare":
    if len(docs) < 2:
        st.info(f"Upload at least 2 PDFs to compare documents. Currently: {len(docs)}.")
    else:
        doc_names = [d["filename"] for d in docs]
        col1, col2 = st.columns(2)
        with col1:
            cmp_doc_a = st.selectbox("Document A", doc_names, key="cmp_doc_a_select")
        with col2:
            non_a_indices = [i for i, name in enumerate(doc_names) if name != cmp_doc_a]
            default_b_index = non_a_indices[0] if non_a_indices else 0
            cmp_doc_b = st.selectbox("Document B", doc_names, key="cmp_doc_b_select", index=default_b_index)

        if cmp_doc_a == cmp_doc_b:
            st.caption("⚠️ Pick two different documents to compare.")

        if st.button("Compare Documents", disabled=(cmp_doc_a == cmp_doc_b)):
            with st.spinner(f"Comparing {cmp_doc_a} vs {cmp_doc_b}..."):
                text_a = " ".join(c.text for c in vs.chunks if c.filename == cmp_doc_a)
                text_b = " ".join(c.text for c in vs.chunks if c.filename == cmp_doc_b)
                try:
                    result = compare_documents(text_a, cmp_doc_a, text_b, cmp_doc_b)
                    st.session_state.comparison_data = {
                        "result": result, "doc_a": cmp_doc_a, "doc_b": cmp_doc_b,
                    }
                    db.log_tool_usage(st.session_state.session_id, "compare")
                except (ValueError, LLMError, DocComparisonError) as e:
                    st.error(str(e))
                    st.session_state.pop("comparison_data", None)

        if st.session_state.get("comparison_data"):
            if st.button("🗑️ Clear comparison"):
                st.session_state.pop("comparison_data", None)
                st.rerun()
            cdata = st.session_state.comparison_data
            result = cdata["result"]
            doc_a, doc_b = cdata["doc_a"], cdata["doc_b"]
            st.markdown(f"**Comparing:** `{doc_a}` vs `{doc_b}`")

            if result["similarities"]:
                st.markdown("#### ✅ Similarities")
                for s in result["similarities"]:
                    st.markdown(f"- {s}")

            if result["differences"]:
                st.markdown("#### 🔀 Differences")
                def _esc(s):
                    return str(s).replace("|", "\\|").replace("\n", " ")
                table_lines = [f"| Aspect | {_esc(doc_a)} | {_esc(doc_b)} |", "|---|---|---|"]
                for d in result["differences"]:
                    table_lines.append(f"| {_esc(d['aspect'])} | {_esc(d['document_a'])} | {_esc(d['document_b'])} |")
                st.markdown("\n".join(table_lines))

            col_a2, col_b2 = st.columns(2)
            with col_a2:
                if result["missing_in_a"]:
                    st.markdown(f"#### 🚫 Missing from `{doc_a}`")
                    for m in result["missing_in_a"]:
                        st.markdown(f"- {m}")
            with col_b2:
                if result["missing_in_b"]:
                    st.markdown(f"#### 🚫 Missing from `{doc_b}`")
                    for m in result["missing_in_b"]:
                        st.markdown(f"- {m}")

            if result["risks"]:
                st.markdown("#### ⚠️ Flagged Risks")
                severity_icon = {"High": "🔴", "Medium": "🟡", "Low": "🟢"}
                for r in result["risks"]:
                    icon = severity_icon.get(r["severity"], "🟡")
                    st.markdown(f"{icon} **{r['severity']}** (applies to: {r['document']}) — {r['description']}")

# ============================== Deep Research ==============================
elif active_tool == "research":
    if not docs:
        st.info("Upload a PDF first to use this tool.")
    else:
        st.caption(
            "Breaks your question into a search plan, retrieves evidence for each "
            "sub-question separately, reasons step by step, then answers with an "
            "honest confidence score — slower than regular chat, but more thorough."
        )
        dr_question = st.text_area(
            "Research question", placeholder="e.g. What are the strengths and weaknesses of this approach?",
            key="deep_research_question", height=80,
        )
        dr_num_subqueries = st.slider("Number of sub-questions to investigate", 2, 6, 4, key="deep_research_num_subq")

        if st.button("🔬 Run Deep Research"):
            with st.spinner("Planning research approach..."):
                try:
                    dr_result = run_deep_research(dr_question, vs, num_subqueries=dr_num_subqueries)
                    st.session_state.deep_research_data = dr_result
                    db.log_tool_usage(st.session_state.session_id, "research")
                except (ValueError, LLMError, DeepResearchError) as e:
                    st.error(str(e))
                    st.session_state.pop("deep_research_data", None)

        if st.session_state.get("deep_research_data"):
            if st.button("🗑️ Clear research"):
                st.session_state.pop("deep_research_data", None)
                st.rerun()
            dr_data = st.session_state.deep_research_data

            st.markdown("#### 🗺️ Search Plan")
            for i, subq in enumerate(dr_data["search_plan"], start=1):
                st.markdown(f"{i}. {subq}")

            if dr_data["reasoning_steps"]:
                with st.expander("🧠 Step-by-step reasoning"):
                    for i, step in enumerate(dr_data["reasoning_steps"], start=1):
                        st.markdown(f"**Step {i}.** {step}")

            st.markdown("#### ✅ Answer")
            st.markdown(dr_data["answer"])

            conf = dr_data["confidence"]
            conf_icon = {"High": "🟢", "Medium": "🟡", "Low": "🔴"}.get(conf["label"], "🟡")
            st.markdown(f"#### {conf_icon} Confidence: {conf['label']} ({conf['score']}/100)")
            if conf["justification"]:
                st.caption(conf["justification"])

            if dr_data["sources"]:
                with st.expander(f"📚 {len(dr_data['sources'])} source(s) consulted"):
                    for s in dr_data["sources"]:
                        subq_tags = ", ".join(s["subqueries"])
                        st.markdown(f"**[Source {s['index']}] {s['filename']}**, page {s['page']} — score {s['score']:.2f}")
                        st.caption(f"Found via: {subq_tags}")
                        st.markdown(f"> {s['text']}")

# ============================== Chat with Image ==============================
elif active_tool == "image":
    if not docs:
        st.info("Upload a PDF first to use this tool.")
    else:
        img_doc = st.selectbox("Choose a document", [d["filename"] for d in docs], key="img_doc_select")
        pdf_path = _uploaded_file_path(st.session_state.session_id, img_doc)

        if not pdf_path.exists():
            st.warning(
                "The original file for this document isn't available anymore (it was "
                "uploaded before this feature was added, or the session data was moved). "
                "Re-upload the PDF to use image analysis on it."
            )
        else:
            try:
                page_count = get_page_count(str(pdf_path))
            except PDFVisionError as e:
                page_count = None
                st.error(str(e))

            if page_count:
                img_page = st.number_input(
                    "Page number", min_value=1, max_value=page_count, value=1, key="img_page_select",
                )
                img_question = st.text_input(
                    "Ask something specific about this page (optional)",
                    placeholder="e.g. What was the peak value shown in the chart?",
                    key="img_question_input",
                )

                if st.button("🔍 Analyze"):
                    with st.spinner("Reading the page..."):
                        try:
                            result_text = describe_page(str(pdf_path), img_page, question=img_question)
                            st.session_state.image_analysis_data = result_text
                            db.log_tool_usage(st.session_state.session_id, "image")
                        except (PDFVisionError, LLMError) as e:
                            st.error(str(e))
                            st.session_state.pop("image_analysis_data", None)

                if st.session_state.get("image_analysis_data"):
                    if st.button("🗑️ Clear analysis"):
                        st.session_state.pop("image_analysis_data", None)
                        st.rerun()
                    st.markdown("#### 🔎 Analysis")
                    st.markdown(st.session_state.image_analysis_data)

# ============================== Translate ==============================
elif active_tool == "translate":
    COMMON_LANGUAGES = [
        "Spanish", "French", "German", "Hindi", "Chinese (Simplified)", "Japanese",
        "Arabic", "Portuguese", "Russian", "Korean", "Italian", "Other (type below)",
    ]
    source_mode = st.radio(
        "What do you want to translate?", ["Paste text", "An uploaded document"],
        key="translate_source_mode",
    )

    text_to_translate = ""
    if source_mode == "Paste text":
        text_to_translate = st.text_area("Text to translate", height=150, key="translate_input_text")
    elif not docs:
        st.info("Upload a document first, or switch to 'Paste text' above.")
    else:
        chosen_doc = st.selectbox("Choose a document", [d["filename"] for d in docs], key="translate_doc_select")
        matching_chunks = [c for c in vs.chunks if c.filename == chosen_doc]
        text_to_translate = " ".join(c.text for c in matching_chunks)

    lang_choice = st.selectbox("Target language", COMMON_LANGUAGES, key="translate_lang_select")
    if lang_choice == "Other (type below)":
        lang_choice = st.text_input("Type the target language", key="translate_lang_custom")

    if st.button("Translate"):
        with st.spinner("Translating..."):
            try:
                result = translate_text(text_to_translate, lang_choice)
                st.session_state.translate_data = result
                db.log_tool_usage(st.session_state.session_id, "translate")
            except (ValueError, LLMError) as e:
                st.error(str(e))
                st.session_state.pop("translate_data", None)

    if st.session_state.get("translate_data"):
        if st.button("🗑️ Clear translation"):
            st.session_state.pop("translate_data", None)
            st.rerun()
        st.markdown("#### 🌍 Translation")
        st.markdown(st.session_state.translate_data)

# ============================== Voice Chat ==============================
elif active_tool == "voice":
    st.caption("Speech-to-text uses Groq's free Whisper API (no extra cost).")

    def _voice_mark_reanswer():
        st.session_state.voice_should_answer = True

    st.session_state.setdefault("voice_audio_widget_id", 0)
    audio_value = st.audio_input(
        "🎤 Record your question", key=f"voice_audio_input_{st.session_state.voice_audio_widget_id}"
    )

    if audio_value is not None:
        import hashlib
        audio_bytes = audio_value.getvalue()
        audio_hash = hashlib.md5(audio_bytes).hexdigest()
        if st.session_state.get("voice_last_audio_hash") != audio_hash:
            with st.spinner("Transcribing..."):
                try:
                    transcript = transcribe_audio(audio_bytes, "recording.wav")
                    st.session_state.voice_transcript = transcript
                    st.session_state.voice_last_audio_hash = audio_hash
                    st.session_state.voice_should_answer = True  # answer automatically, no Send click needed
                except LLMError as e:
                    st.error(str(e))
                    st.session_state.pop("voice_transcript", None)

    if st.session_state.get("voice_transcript"):
        if st.button("🗑️ Clear recording"):
            # Bump the audio widget's key so it resets to empty too — clearing
            # only our own session_state variables leaves the recorded clip
            # still showing in the widget itself, since Streamlit widgets keep
            # their own state tied to their key independent of anything else.
            st.session_state.voice_audio_widget_id += 1
            for k in (
                "voice_transcript", "voice_last_audio_hash", "voice_edited_question",
                "voice_should_answer", "voice_answer_data",
            ):
                st.session_state.pop(k, None)
            st.rerun()

        edited_question = st.text_area(
            "Transcribed question (edit to ask something different — updates automatically)",
            value=st.session_state.voice_transcript, key="voice_edited_question",
            on_change=_voice_mark_reanswer,
        )

        if st.session_state.get("voice_should_answer"):
            st.session_state.voice_should_answer = False
            if not vs.has_documents():
                st.warning("Please upload at least one PDF first.")
            elif not edited_question.strip():
                st.warning("The transcribed question is empty — try recording again.")
            else:
                with st.spinner("Searching documents..."):
                    retrieved = vs.search(edited_question, top_k=top_k)
                try:
                    answer = st.write_stream(generate_answer_stream(edited_question, retrieved))
                except LLMError as e:
                    answer = f"⚠️ {e}"
                    st.markdown(answer)

                sources = [
                    {"filename": r.filename, "page": r.page, "score": r.score, "text": r.text}
                    for r in retrieved
                ]
                if sources:
                    with st.expander(f"📎 {len(sources)} source excerpt(s)"):
                        for s in sources:
                            st.markdown(f"**{s['filename']}**, page {s['page']} — score {s['score']:.2f}")
                            st.markdown(f"> {s['text']}")

                db.add_chat_entry(st.session_state.session_id, edited_question, answer, sources)
                st.session_state.chat_history.append(
                    {"question": edited_question, "answer": answer, "sources": sources}
                )
                db.log_tool_usage(st.session_state.session_id, "voice")
                st.session_state.voice_answer_data = {"answer": answer, "sources": sources}
        elif st.session_state.get("voice_answer_data"):
            # Redisplay the last answer on reruns that weren't triggered by a
            # new question, so it doesn't disappear from view.
            prev = st.session_state.voice_answer_data
            st.markdown(prev["answer"])
            if prev["sources"]:
                with st.expander(f"📎 {len(prev['sources'])} source excerpt(s)"):
                    for s in prev["sources"]:
                        st.markdown(f"**{s['filename']}**, page {s['page']} — score {s['score']:.2f}")
                        st.markdown(f"> {s['text']}")