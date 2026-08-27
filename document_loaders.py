"""
document_loaders.py
Extracts text content from various file formats as a list of "pages" — one
string per logical unit (PDF page, PPTX slide, XLSX sheet, or a single block
for flat text formats). This keeps the existing page-aware citation system
(filename + page number) working uniformly across formats.

Images are handled specially: since there's no text to "extract" from a photo
or scanned page, the file is sent to a vision model to generate a description
(including any readable text via OCR), and that description becomes the
indexed "page" content.
"""

import csv
import io
from pathlib import Path

import docx  # python-docx
import openpyxl
from pptx import Presentation  # python-pptx

from llm_client import analyze_image, LLMError
from pdf_vision import resize_image_bytes

# Extensions handled as plain UTF-8 text (source code, markup, config, etc.)
PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".go", ".rb", ".php", ".rs", ".swift", ".kt", ".sql",
    ".html", ".css", ".json", ".yaml", ".yml", ".xml", ".sh", ".bat",
}

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".tif"}

# All extensions this module claims to support, for use in file_uploader's `type=` list.
# Note: ".pdf" is included here for that purpose, but load_document() below does NOT
# handle it — PDF extraction stays in rag_engine.py (via pdfplumber, already tested)
# to avoid a circular import (rag_engine imports this module, not the other way around).
# VectorStore.add_document() special-cases ".pdf" before falling through to this module.
SUPPORTED_EXTENSIONS = (
    {".pdf", ".docx", ".pptx", ".xlsx", ".csv"} | PLAIN_TEXT_EXTENSIONS | IMAGE_EXTENSIONS
)


class DocumentLoadError(Exception):
    """Raised when a file can't be parsed into text content."""
    pass


def _load_docx(path: str) -> list:
    document = docx.Document(path)
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    # Tables often carry real content (e.g. resumes, spec sheets) — include them too
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    full_text = "\n".join(paragraphs)
    return [full_text] if full_text.strip() else []


def _load_pptx(path: str) -> list:
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        pieces = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text.strip():
                    pieces.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        pieces.append(row_text)
        slides_text.append("\n".join(pieces))
    return slides_text


def _load_xlsx(path: str) -> list:
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    sheets_text = []
    for sheet in wb.worksheets:
        lines = [f"Sheet: {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
        sheets_text.append("\n".join(lines))
    wb.close()
    return sheets_text


def _load_csv(path: str) -> list:
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        lines = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
    return ["\n".join(lines)] if lines else []


def _load_plain_text(path: str) -> list:
    with open(path, encoding="utf-8", errors="replace") as f:
        content = f.read()
    return [content] if content.strip() else []


def _load_image(path: str, ext: str) -> list:
    """Describe an uploaded image using the vision model so its content becomes
    searchable text. Raises LLMError if no Groq key is configured (the only
    vision-capable provider in this app)."""
    with open(path, "rb") as f:
        raw_bytes = f.read()
    import base64
    # resize_image_bytes always re-encodes to PNG, so the MIME type sent to the
    # API must be image/png regardless of the original upload's extension
    raw_bytes = resize_image_bytes(raw_bytes)
    b64 = base64.b64encode(raw_bytes).decode("utf-8")
    mime_type = "image/png"

    prompt = (
        "Describe this image thoroughly for a search index: transcribe any visible text "
        "verbatim (OCR), describe any charts/graphs/tables with their actual data, and "
        "describe any diagrams, photos, or handwriting. Be comprehensive — this is the "
        "only representation of this image that will be searchable."
    )
    description = analyze_image(b64, mime_type=mime_type, prompt=prompt, max_tokens=1500)
    return [description]


def load_document(file_path: str, filename: str) -> list:
    """Load a file and return a list of page/section text strings.

    Dispatches based on file extension. Does NOT handle ".pdf" — that's handled
    directly by rag_engine.py's VectorStore.add_document(), which special-cases
    it before falling through here (see the SUPPORTED_EXTENSIONS comment above).

    Raises:
        DocumentLoadError: unsupported extension (including ".pdf", if called
            directly with one — callers should route PDFs elsewhere).
        LLMError: an image needs analysis but no vision provider is configured.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        raise DocumentLoadError(
            "load_document() doesn't handle PDFs directly — use "
            "VectorStore.add_document(), which routes PDFs to rag_engine's own "
            "pdfplumber-based extraction."
        )
    if ext == ".docx":
        return _load_docx(file_path)
    if ext == ".pptx":
        return _load_pptx(file_path)
    if ext == ".xlsx":
        return _load_xlsx(file_path)
    if ext == ".csv":
        return _load_csv(file_path)
    if ext in PLAIN_TEXT_EXTENSIONS:
        return _load_plain_text(file_path)
    if ext in IMAGE_EXTENSIONS:
        return _load_image(file_path, ext)

    raise DocumentLoadError(
        f"Unsupported file type: '{ext}'. Supported: PDF, DOCX, PPTX, XLSX, CSV, "
        f"images (jpg/png/webp/bmp/tiff), and plain text/code files."
    )